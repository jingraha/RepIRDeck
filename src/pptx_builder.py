from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .config import (
    ACTIVE_SLIDE_IDS,
    COLORS,
    FONT_MONO,
    FONT_SANS,
    PPTX_PATH,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
)
from .model import ModelBundle
from .pptx_helpers import (
    add_bullet_list,
    add_card_title,
    add_chart_legend,
    add_chevron,
    add_circle,
    add_horizontal_bars,
    add_line,
    add_matrix_dot,
    add_metric_card,
    add_native_line_chart,
    add_rect,
    add_replit_mark,
    add_rich_text,
    add_stacked_columns,
    add_tag,
    add_text,
    rgb,
    set_background,
)
from .sources import Source, format_slide_sources


class DeckBuilder:
    def __init__(
        self,
        model: ModelBundle,
        sources: dict[str, Source],
        slide_content: dict[str, Any],
    ) -> None:
        self.model = model
        self.sources = sources
        self.content = slide_content
        self.source_map = slide_content["source_map"]
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.blank_layout = self.prs.slide_layouts[6]
        self.display_slide_number = 0

    def _new_slide(
        self,
        number: int,
        *,
        dark: bool = False,
        show_header: bool = True,
    ) -> Any:
        slide = self.prs.slides.add_slide(self.blank_layout)
        self.display_slide_number += 1
        set_background(slide, "ink" if dark else "cream")
        if show_header:
            slide_copy = self.content["slides"][number]
            title_size = 24 if len(slide_copy["title"]) <= 58 else 19.5
            add_text(
                slide,
                slide_copy["title"],
                0.65,
                0.31,
                10.95,
                0.72,
                font_size=title_size,
                color="white" if dark else "ink",
                bold=True,
            )
            add_text(
                slide,
                slide_copy["subtitle"],
                0.65,
                1.05,
                11.1,
                0.32,
                font_size=9.5,
                color="gray_300" if dark else "gray_700",
            )
            add_replit_mark(
                slide,
                12.0,
                0.45,
                scale=0.43,
                wordmark=False,
                color="orange",
            )
            add_line(
                slide,
                0.65,
                1.39,
                12.68,
                1.39,
                color="charcoal" if dark else "gray_200",
                width=0.8,
            )
        self._footer(
            slide,
            number,
            self.display_slide_number,
            dark=dark,
        )
        return slide

    def _footer(
        self,
        slide: Any,
        content_number: int,
        display_number: int,
        *,
        dark: bool = False,
    ) -> None:
        source_text = format_slide_sources(
            content_number,
            self.sources,
            self.source_map,
        )
        add_text(
            slide,
            source_text,
            0.65,
            7.16,
            8.4,
            0.16,
            font_size=5.4,
            color="gray_500" if dark else "gray_700",
        )
        add_text(
            slide,
            "Illustrative | Public information + analyst estimates",
            9.1,
            7.16,
            3.05,
            0.16,
            font_size=5.4,
            color="gray_500" if dark else "gray_700",
            align=PP_ALIGN.RIGHT,
        )
        add_text(
            slide,
            f"{display_number:02d}",
            12.35,
            7.12,
            0.34,
            0.20,
            font_size=7,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
            align=PP_ALIGN.RIGHT,
        )

    def build(self, output_path: Path = PPTX_PATH) -> Path:
        slide_methods = [
            self._slide_01_cover,
            self._slide_02_transaction,
            self._slide_03_thesis,
            self._slide_04_at_a_glance,
            self._slide_05_vision,
            self._slide_06_problem,
            self._slide_07_platform,
            self._slide_08_agent_workflow,
            self._slide_09_customer_value,
            self._slide_10_why_now,
            self._slide_11_market,
            self._slide_12_category_leadership,
            self._slide_19_gtm,
            self._slide_21_competition,
            self._slide_23_historical_financials,
            self._slide_24_operating_plan,
            self._slide_15_retention,
            self._slide_18_unit_economics,
            self._slide_25_profitability,
            self._slide_26_use_of_proceeds,
            self._slide_27_milestones,
            self._slide_28_leadership,
            self._slide_29_risks,
            self._slide_30_closing,
        ]
        if len(slide_methods) != len(ACTIVE_SLIDE_IDS):
            raise ValueError("Active slide list does not match the slide methods")
        for method in slide_methods:
            method()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.core_properties.title = "Replit Series E Investor Presentation"
        self.prs.core_properties.subject = (
            "Illustrative late-stage financing materials based on public information"
        )
        self.prs.core_properties.author = "Replit Series E materials generator"
        self.prs.save(output_path)
        return output_path

    def _slide_01_cover(self) -> None:
        slide = self._new_slide(1, dark=True, show_header=False)
        tx = self.model.transaction
        add_replit_mark(
            slide,
            0.75,
            0.62,
            scale=0.78,
            wordmark=True,
            color="orange",
            wordmark_color="white",
        )
        add_tag(
            slide,
            "Series E",
            11.45,
            0.68,
            fill="orange",
            color="white",
            width=1.05,
        )
        add_text(
            slide,
            "Build and run software\nfrom one prompt",
            0.75,
            1.68,
            9.5,
            2.05,
            font_size=42,
            color="white",
            bold=True,
            line_spacing=0.88,
        )
        add_text(
            slide,
            "Series E financing | $1B raise | $20B pre-money",
            0.78,
            3.93,
            8.7,
            0.42,
            font_size=15,
            color="gray_300",
        )
        cards = [
            (f"${tx['total_raise'] / 1000:.1f}B", "Total raise", "orange"),
            (f"${tx['pre_money'] / 1000:.0f}B", "Pre-money", "blue"),
            (f"${tx['primary']:.0f}M", "Primary", "coral"),
        ]
        add_tag(
            slide,
            "Illustrative financing case",
            0.78,
            4.65,
            fill="charcoal",
            color="gray_300",
            width=1.78,
        )
        for index, (value, label, accent) in enumerate(cards):
            add_metric_card(
                slide,
                0.78 + index * 2.35,
                5.05,
                2.05,
                1.18,
                value,
                label,
                accent=accent,
                fill="charcoal",
                dark=True,
                value_size=23,
            )
        add_text(
            slide,
            "CONFIDENTIAL",
            10.65,
            5.20,
            1.85,
            0.30,
            font_size=8,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
            align=PP_ALIGN.RIGHT,
        )
        add_text(
            slide,
            "September 2026",
            10.65,
            5.62,
            1.85,
            0.30,
            font_size=10,
            color="white",
            align=PP_ALIGN.RIGHT,
        )

    def _slide_02_transaction(self) -> None:
        slide = self._new_slide(2)
        tx = self.model.transaction
        cards = [
            (f"${tx['total_raise'] / 1000:.1f}B", "Total financing", "orange", "Est."),
            (f"${tx['primary']:.0f}M", "Primary capital", "blue", "Est."),
            (f"${tx['secondary']:.0f}M", "Secondary liquidity", "coral", "Est."),
            (
                f"{tx['total_buyer_ownership']:.1%}",
                f"Total buyer ownership\n{tx['primary_ownership']:.1%} primary dilution",
                "charcoal",
                "Derived",
            ),
        ]
        for index, (value, label, accent, evidence) in enumerate(cards):
            add_metric_card(
                slide,
                0.65 + (index % 2) * 2.7,
                1.70 + (index // 2) * 1.55,
                2.45,
                1.30,
                value,
                label,
                accent=accent,
                evidence=evidence,
                value_size=24,
            )

        add_text(
            slide,
            "VALUATION CONTEXT",
            6.35,
            1.72,
            2.0,
            0.26,
            font_size=8,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
        )
        bars = [
            ("Mar-26 Series D", tx["prior_round_valuation"] / 1000, "gray_500"),
            ("Series E pre-money", tx["pre_money"] / 1000, "orange"),
            ("Series E post-money", tx["post_money"] / 1000, "blue"),
        ]
        add_horizontal_bars(
            slide,
            6.35,
            2.12,
            5.95,
            [item[0] for item in bars],
            [item[1] for item in bars],
            [item[2] for item in bars],
            max_value=22,
            value_format="${:,.1f}B",
            row_height=0.80,
        )
        add_text(
            slide,
            f"{tx['valuation_step_up']:.1f}x",
            6.35,
            4.77,
            1.1,
            0.48,
            font_size=24,
            color="orange",
            bold=True,
        )
        add_text(
            slide,
            "step-up vs. March 2026 round",
            7.35,
            4.83,
            2.4,
            0.30,
            font_size=9,
            color="gray_700",
        )
        add_text(
            slide,
            f"{tx['pre_money'] / self.model.financials.loc[2027, 'revenue']:.1f}x",
            9.55,
            4.77,
            1.15,
            0.48,
            font_size=24,
            color="blue",
            bold=True,
        )
        add_text(
            slide,
            "pre-money / 2027E revenue",
            10.55,
            4.83,
            1.77,
            0.30,
            font_size=8.3,
            color="gray_700",
        )
        add_rect(slide, 0.65, 5.58, 11.98, 0.96, fill="ink", radius=True)
        add_rich_text(
            slide,
            [
                {"text": "Primary: ", "font_size": 11, "bold": True, "color": "white"},
                {
                    "text": "funds models, users, enterprise sales, international growth, and talent",
                    "font_size": 9.5,
                    "color": "gray_300",
                },
                {"text": "   |   Secondary: ", "font_size": 9.5, "bold": True, "color": "white"},
                {
                    "text": "employee and early-investor liquidity",
                    "font_size": 9.5,
                    "color": "gray_300",
                },
            ],
            0.95,
            5.84,
            11.2,
            0.38,
            valign=MSO_ANCHOR.MIDDLE,
        )

    def _slide_03_thesis(self) -> None:
        slide = self._new_slide(3)
        cards = [
            (
                "01",
                "Distribution",
                "50M+ users give Replit a direct path to teams and enterprise accounts.",
                "orange",
            ),
            (
                "02",
                "Full-stack platform",
                "Agent, workspace, collaboration, deployment, and runtime are in one product.",
                "blue",
            ),
            (
                "03",
                "Enterprise demand",
                "Users from 85% of the Fortune 500 show demand inside large companies.",
                "coral",
            ),
            (
                "04",
                "Improving margins",
                "Lower model cost and more enterprise revenue lift gross margin.",
                "green",
            ),
            (
                "05",
                "Several revenue sources",
                "Replit earns revenue from individuals, teams, enterprises, agents, and running apps.",
                "charcoal",
            ),
        ]
        positions = [
            (0.65, 1.72, 3.78, 1.72),
            (4.78, 1.72, 3.78, 1.72),
            (8.90, 1.72, 3.78, 1.72),
            (2.72, 3.78, 3.78, 1.72),
            (6.84, 3.78, 3.78, 1.72),
        ]
        for card, position in zip(cards, positions, strict=True):
            index, title, body, accent = card
            add_card_title(
                slide,
                index,
                title,
                body,
                *position,
                accent=accent,
            )
        add_rect(slide, 0.65, 5.87, 12.03, 0.64, fill="orange", radius=True)
        add_text(
            slide,
            "AI changes software creation. The key question is who owns the full workflow.",
            0.95,
            6.02,
            11.4,
            0.30,
            font_size=11,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def _slide_04_at_a_glance(self) -> None:
        slide = self._new_slide(4)
        f = self.model.financials.loc[2026]
        anchors = self.model.assumptions["public_anchors"]
        company = self.model.assumptions["company_estimates"]
        metrics = [
            (f"{f['registered_users_m']:.0f}M+", "Registered users", "PUBLIC", "orange"),
            (f"{anchors['fortune_500_share']:.0%}", "Fortune 500 presence", "PUBLIC", "blue"),
            (f"${f['exit_arr'] / 1000:.1f}B", "2026 exit ARR goal", "MT", "coral"),
            (f"${f['revenue']:,.0f}M", "2026 revenue", "EST.", "green"),
            (f"{f['enterprise_customers']:,.0f}", "Enterprise customers", "EST.", "charcoal"),
            (f"{company['countries_reached']:.0f}+", "Countries reached", "EST.", "gray_500"),
        ]
        for index, (value, label, evidence, accent) in enumerate(metrics):
            column = index % 3
            row = index // 3
            add_metric_card(
                slide,
                0.65 + column * 4.05,
                1.72 + row * 1.62,
                3.68,
                1.36,
                value,
                label,
                accent=accent,
                evidence=evidence,
                value_size=26,
            )
        add_rect(slide, 0.65, 5.17, 12.03, 1.33, fill="ink", radius=True)
        add_text(
            slide,
            "One product from idea to running app",
            0.98,
            5.34,
            4.5,
            0.56,
            font_size=13.5,
            color="white",
            bold=True,
        )
        add_text(
            slide,
            "Distribution",
            6.05,
            5.34,
            1.25,
            0.22,
            font_size=7.5,
            color="gray_300",
            bold=True,
            all_caps=True,
        )
        add_text(
            slide,
            "Agent",
            7.55,
            5.34,
            1.25,
            0.22,
            font_size=7.5,
            color="gray_300",
            bold=True,
            all_caps=True,
        )
        add_text(
            slide,
            "Enterprise",
            9.05,
            5.34,
            1.25,
            0.22,
            font_size=7.5,
            color="gray_300",
            bold=True,
            all_caps=True,
        )
        add_text(
            slide,
            "Runtime",
            10.85,
            5.34,
            1.25,
            0.22,
            font_size=7.5,
            color="gray_300",
            bold=True,
            all_caps=True,
        )
        for index, color in enumerate(("orange", "blue", "coral", "green")):
            add_rect(slide, 6.05 + index * 1.52, 5.77, 1.18, 0.18, fill=color, radius=True)
            if index < 3:
                add_chevron(slide, 7.30 + index * 1.52, 5.73, 0.19, 0.26, fill="gray_500")
        add_text(
            slide,
            "One revenue loop",
            6.05,
            6.08,
            6.0,
            0.24,
            font_size=8,
            color="gray_300",
            align=PP_ALIGN.CENTER,
        )

    def _slide_05_vision(self) -> None:
        slide = self._new_slide(5)
        users = self.model.financials.loc[2026, "registered_users_m"]
        add_rect(slide, 0.65, 1.72, 5.10, 4.78, fill="ink", radius=True)
        add_text(
            slide,
            f"{users:.0f}M+",
            0.98,
            2.18,
            3.6,
            0.80,
            font_size=48,
            color="orange",
            bold=True,
        )
        add_text(
            slide,
            "creators today",
            1.02,
            3.01,
            3.5,
            0.42,
            font_size=17,
            color="white",
            bold=True,
        )
        add_line(slide, 1.02, 3.69, 5.12, 3.69, color="charcoal", width=1)
        add_text(
            slide,
            "Replit's ambition",
            1.02,
            4.06,
            2.0,
            0.25,
            font_size=8,
            color="gray_300",
            bold=True,
            all_caps=True,
            font_name=FONT_MONO,
        )
        add_text(
            slide,
            "Billions turn intent\ninto working software",
            1.02,
            4.47,
            3.95,
            1.12,
            font_size=25,
            color="white",
            bold=True,
        )

        personas = [
            ("STUDENT", "Learn by building", "orange"),
            ("FOUNDER", "Launch without a team", "blue"),
            ("EMPLOYEE", "Automate the backlog", "coral"),
            ("ENTERPRISE", "Govern creation at scale", "green"),
        ]
        for index, (persona, promise, accent) in enumerate(personas):
            y = 1.72 + index * 1.18
            add_rect(slide, 6.13, y, 6.55, 0.94, fill="white", line="gray_200", radius=True)
            add_circle(slide, 6.40, y + 0.20, 0.50, fill=accent)
            add_text(
                slide,
                str(index + 1),
                6.40,
                y + 0.20,
                0.50,
                0.50,
                font_size=10,
                color="white",
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                font_name=FONT_MONO,
            )
            add_text(
                slide,
                persona,
                7.15,
                y + 0.17,
                1.45,
                0.22,
                font_size=7.5,
                color=accent,
                bold=True,
                all_caps=True,
                font_name=FONT_MONO,
            )
            add_text(
                slide,
                promise,
                7.15,
                y + 0.46,
                4.7,
                0.27,
                font_size=12,
                color="ink",
                bold=True,
            )
        add_text(
            slide,
            "Intent",
            6.30,
            6.38,
            1.0,
            0.25,
            font_size=9,
            color="gray_700",
            bold=True,
        )
        add_chevron(slide, 7.22, 6.36, 0.34, 0.28, fill="orange")
        add_text(
            slide,
            "Agentic creation",
            7.72,
            6.38,
            1.7,
            0.25,
            font_size=9,
            color="gray_700",
            bold=True,
        )
        add_chevron(slide, 9.42, 6.36, 0.34, 0.28, fill="blue")
        add_text(
            slide,
            "Deployed software",
            9.92,
            6.38,
            1.9,
            0.25,
            font_size=9,
            color="gray_700",
            bold=True,
        )

    def _slide_06_problem(self) -> None:
        slide = self._new_slide(6)
        benchmarks = self.model.assumptions["problem_benchmarks"]
        steps = [
            ("SPEC", "Translate need", "orange"),
            ("DESIGN", "Prototype", "coral"),
            ("CODE", "Build", "blue"),
            ("TEST", "Validate", "green"),
            ("SECURE", "Review", "charcoal"),
            ("DEPLOY", "Operate", "orange"),
        ]
        for index, (stage, action, color) in enumerate(steps):
            x = 0.65 + index * 2.02
            add_rect(slide, x, 1.92, 1.67, 1.35, fill="white", line="gray_200", radius=True)
            add_text(
                slide,
                stage,
                x + 0.18,
                2.13,
                1.25,
                0.22,
                font_size=7,
                color=color,
                bold=True,
                font_name=FONT_MONO,
                align=PP_ALIGN.CENTER,
            )
            add_text(
                slide,
                action,
                x + 0.12,
                2.55,
                1.42,
                0.30,
                font_size=11,
                color="ink",
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            if index < len(steps) - 1:
                add_chevron(slide, x + 1.75, 2.40, 0.20, 0.34, fill="gray_300")
        add_text(
            slide,
            "Every handoff introduces latency, translation loss, and another queue.",
            0.65,
            3.58,
            12.03,
            0.40,
            font_size=15,
            color="ink",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        estimates = [
            (f"{benchmarks['specialized_tools']}+", "specialized tools", "orange"),
            (
                f"{benchmarks['handoffs_low']} to {benchmarks['handoffs_high']}",
                "cross-functional handoffs",
                "blue",
            ),
            (
                f"{benchmarks['cycle_weeks_low']} to {benchmarks['cycle_weeks_high']} wks",
                "idea-to-production cycle",
                "coral",
            ),
            (
                f"<${benchmarks['agent_test_cost_usd']}",
                "cost to test an idea with agents",
                "green",
            ),
        ]
        for index, (value, label, color) in enumerate(estimates):
            add_metric_card(
                slide,
                0.65 + index * 3.05,
                4.38,
                2.72,
                1.40,
                value,
                label,
                accent=color,
                evidence="EST.",
                value_size=22,
            )
        add_rect(slide, 0.65, 6.03, 12.03, 0.48, fill="gray_100", radius=True)
        add_text(
            slide,
            "The scarce resource is coordinated work from an idea to a secure, running product.",
            0.90,
            6.14,
            11.55,
            0.25,
            font_size=9.5,
            color="gray_700",
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def _slide_07_platform(self) -> None:
        slide = self._new_slide(7)
        layers = [
            ("Intent layer", "Natural language • Canvas • Team tasks", "orange"),
            ("Agent layer", "Plan • Build • Test • Repair • Security", "blue"),
            ("Workspace layer", "Code • Data • Collaboration • Versioning", "coral"),
            ("Runtime layer", "Deploy • Compute • Databases • Observability", "green"),
            ("Control layer", "Identity • Governance • Policy • Audit", "charcoal"),
        ]
        for index, (layer, details, color) in enumerate(layers):
            y = 1.72 + index * 0.88
            width = 8.25 - index * 0.28
            x = 0.65 + index * 0.14
            add_rect(slide, x, y, width, 0.66, fill=color, radius=True)
            add_text(
                slide,
                layer,
                x + 0.24,
                y + 0.16,
                1.65,
                0.24,
                font_size=10,
                color="white",
                bold=True,
            )
            add_text(
                slide,
                details,
                x + 2.00,
                y + 0.16,
                width - 2.30,
                0.24,
                font_size=9,
                color="white",
                align=PP_ALIGN.RIGHT,
            )
        add_rect(slide, 9.32, 1.72, 3.36, 4.18, fill="ink", radius=True)
        add_text(
            slide,
            "ONE CLOSED LOOP",
            9.64,
            2.02,
            2.75,
            0.25,
            font_size=8,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
            align=PP_ALIGN.CENTER,
        )
        loop = [
            ("BUILD", "orange"),
            ("LEARN", "blue"),
            ("DEPLOY", "coral"),
            ("IMPROVE", "green"),
        ]
        center_x, center_y = 10.99, 3.75
        node_diameter = 0.72
        positions = [
            (center_x - 0.36, center_y - 1.20),
            (center_x + 0.69, center_y - 0.36),
            (center_x - 0.36, center_y + 0.73),
            (center_x - 1.41, center_y - 0.36),
        ]
        for x2, y2 in (
            (center_x, center_y - 0.84),
            (center_x + 1.05, center_y),
            (center_x, center_y + 1.09),
            (center_x - 1.05, center_y),
        ):
            add_line(
                slide,
                center_x,
                center_y,
                x2,
                y2,
                color="gray_500",
                width=1,
            )
        for (label, color), (x, y) in zip(loop, positions, strict=True):
            add_circle(slide, x, y, node_diameter, fill=color)
            add_text(
                slide,
                label,
                x,
                y,
                node_diameter,
                node_diameter,
                font_size=6.3,
                color="white",
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                font_name=FONT_MONO,
            )
        add_circle(slide, center_x - 0.37, center_y - 0.37, 0.74, fill="white")
        add_text(
            slide,
            "DATA",
            center_x - 0.37,
            center_y - 0.37,
            0.74,
            0.74,
            font_size=7.5,
            color="ink",
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            font_name=FONT_MONO,
        )
        add_rect(slide, 0.65, 6.15, 12.03, 0.37, fill="white", line="gray_200", radius=True)
        add_text(
            slide,
            "The product records the build workflow and the running workload. This data improves the next build.",
            0.88,
            6.23,
            11.60,
            0.20,
            font_size=8.7,
            color="gray_700",
            align=PP_ALIGN.CENTER,
        )

    def _slide_08_agent_workflow(self) -> None:
        slide = self._new_slide(8)
        agent_speed = self.model.assumptions["public_anchors"]["agent_4_speed_multiple"]
        add_rect(slide, 0.65, 1.72, 2.15, 4.80, fill="ink", radius=True)
        add_text(
            slide,
            "HUMAN",
            0.95,
            2.04,
            1.55,
            0.25,
            font_size=8,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
            align=PP_ALIGN.CENTER,
        )
        add_circle(slide, 1.28, 2.58, 0.90, fill="orange")
        add_text(
            slide,
            "INTENT",
            1.28,
            2.58,
            0.90,
            0.90,
            font_size=9,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            font_name=FONT_MONO,
        )
        add_bullet_list(
            slide,
            ["Set direction", "Review choices", "Refine outcome"],
            0.98,
            3.85,
            1.55,
            font_size=8,
            color="gray_300",
            row_height=0.48,
        )
        add_chevron(slide, 2.96, 3.72, 0.42, 0.50, fill="orange")

        add_text(
            slide,
            "PARALLEL AGENTS",
            3.65,
            1.78,
            3.0,
            0.25,
            font_size=8,
            color="gray_700",
            bold=True,
            font_name=FONT_MONO,
        )
        agent_cards = [
            ("PLAN", "Decompose requirements", "orange"),
            ("BUILD", "Create product + data", "blue"),
            ("TEST", "Run and repair", "coral"),
            ("SECURE", "Find and patch risk", "green"),
        ]
        for index, (label, body, color) in enumerate(agent_cards):
            row = index % 2
            column = index // 2
            x = 3.55 + column * 2.62
            y = 2.18 + row * 1.63
            add_rect(slide, x, y, 2.27, 1.32, fill="white", line="gray_200", radius=True)
            add_tag(slide, label, x + 0.18, y + 0.15, fill=color, color="white", width=0.78)
            add_text(
                slide,
                body,
                x + 0.18,
                y + 0.61,
                1.90,
                0.44,
                font_size=10,
                color="ink",
                bold=True,
            )
        add_chevron(slide, 8.90, 3.72, 0.42, 0.50, fill="blue")

        add_rect(slide, 9.57, 1.72, 3.11, 4.80, fill="white", line="gray_200", radius=True)
        add_text(
            slide,
            "OUTCOME",
            9.88,
            2.04,
            2.50,
            0.25,
            font_size=8,
            color="blue",
            bold=True,
            font_name=FONT_MONO,
            align=PP_ALIGN.CENTER,
        )
        add_rect(slide, 10.02, 2.58, 2.22, 1.34, fill="blue", radius=True)
        add_text(
            slide,
            "RUNNING\nSOFTWARE",
            10.02,
            2.58,
            2.22,
            1.34,
            font_size=17,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            "Agent 4",
            10.22,
            4.34,
            1.82,
            0.35,
            font_size=22,
            color="orange",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            f"{agent_speed:.0f}x faster than Agent 3",
            9.90,
            4.85,
            2.46,
            0.35,
            font_size=9,
            color="gray_700",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_tag(slide, "PUBLIC", 10.67, 5.40, fill="gray_100", color="gray_700", width=0.92)

    def _slide_09_customer_value(self) -> None:
        slide = self._new_slide(9)
        benchmarks = self.model.assumptions["problem_benchmarks"]
        ukg_gain = self.model.assumptions["public_anchors"][
            "ukg_feedback_capacity_multiple"
        ]
        add_text(
            slide,
            "TRADITIONAL",
            0.65,
            1.78,
            1.50,
            0.22,
            font_size=8,
            color="gray_700",
            bold=True,
            font_name=FONT_MONO,
        )
        add_rect(slide, 0.65, 2.18, 8.15, 0.56, fill="gray_200", radius=True)
        phases = [
            ("Queue", 1.45, "gray_500"),
            ("Specify", 1.35, "charcoal"),
            ("Build", 2.10, "blue"),
            ("Review", 1.35, "coral"),
            ("Deploy", 1.90, "green"),
        ]
        cursor = 0.65
        total = sum(item[1] for item in phases)
        for label, width, color in phases:
            scaled = 8.15 * width / total
            add_rect(slide, cursor, 2.18, scaled, 0.56, fill=color, radius=True)
            add_text(
                slide,
                label,
                cursor,
                2.32,
                scaled,
                0.20,
                font_size=7,
                color="white",
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            cursor += scaled
        add_text(
            slide,
            f"{benchmarks['cycle_weeks_low']} to {benchmarks['cycle_weeks_high']} weeks",
            9.15,
            2.25,
            1.65,
            0.32,
            font_size=15,
            color="ink",
            bold=True,
        )
        add_tag(slide, "EST.", 11.45, 2.29, width=0.62)

        add_text(
            slide,
            "WITH REPLIT",
            0.65,
            3.20,
            1.50,
            0.22,
            font_size=8,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
        )
        add_rect(slide, 0.65, 3.59, 2.25, 0.56, fill="orange", radius=True)
        add_text(
            slide,
            "Prompt → review → deploy",
            0.65,
            3.73,
            2.25,
            0.20,
            font_size=7.5,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "hours to days",
            3.27,
            3.66,
            1.75,
            0.32,
            font_size=15,
            color="orange",
            bold=True,
        )
        add_tag(slide, "EST.", 5.04, 3.70, width=0.62)

        add_rect(slide, 0.65, 4.67, 5.75, 1.44, fill="white", line="gray_200", radius=True)
        add_tag(
            slide,
            "EXAMPLE CUSTOMER: UKG",
            0.95,
            4.82,
            fill="orange",
            color="white",
            width=1.55,
        )
        add_tag(slide, "PUBLIC", 5.18, 4.82, width=0.92)
        add_text(
            slide,
            f"{ukg_gain:.0%}",
            0.96,
            5.20,
            1.60,
            0.55,
            font_size=31,
            color="orange",
            bold=True,
        )
        add_text(
            slide,
            "more customer feedback before engineering starts",
            2.62,
            5.22,
            3.40,
            0.42,
            font_size=10,
            color="ink",
            bold=True,
        )

        add_rect(slide, 6.73, 4.67, 5.95, 1.44, fill="ink", radius=True)
        add_text(
            slide,
            "“Show, not tell.”",
            7.05,
            4.93,
            2.30,
            0.36,
            font_size=18,
            color="white",
            bold=True,
        )
        add_text(
            slide,
            "Product leaders turn rough concepts into functional prototypes before traditional handoffs begin.",
            7.05,
            5.43,
            5.10,
            0.42,
            font_size=9,
            color="gray_300",
        )

    def _slide_10_why_now(self) -> None:
        slide = self._new_slide(10)
        center_x, center_y = 5.53, 2.55
        add_circle(slide, center_x, center_y, 2.20, fill="ink")
        add_text(
            slide,
            "WHY\nNOW",
            center_x,
            center_y,
            2.20,
            2.20,
            font_size=17,
            color="white",
            bold=True,
            font_name=FONT_MONO,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        catalysts = [
            ("01", "Capability", "Agents plan, test, and repair across full applications.", "orange", (0.65, 1.75)),
            ("02", "Economics", "Inference and routing costs decline as quality rises.", "blue", (8.87, 1.75)),
            ("03", "Labor", "Software demand grows faster than engineering capacity.", "coral", (0.65, 4.55)),
            ("04", "Distribution", "Cloud creation reaches users outside traditional IDEs.", "green", (8.87, 4.55)),
            ("05", "Governance", "Identity, security, and policy make enterprise rollout viable.", "charcoal", (4.68, 5.25)),
        ]
        for index, title, body, color, (x, y) in catalysts:
            add_card_title(slide, index, title, body, x, y, 3.82, 1.46, accent=color)
        connectors = [
            (4.47, 2.48, 5.53, 3.16),
            (8.87, 2.48, 7.73, 3.16),
            (4.47, 5.17, 5.69, 4.55),
            (8.87, 5.17, 7.58, 4.55),
            (6.59, 5.25, 6.63, 4.75),
        ]
        for x1, y1, x2, y2 in connectors:
            add_line(slide, x1, y1, x2, y2, color="gray_300", width=1.2)

    def _slide_11_market(self) -> None:
        slide = self._new_slide(11)
        market = self.model.market
        positions = [
            (0.65, 1.72),
            (4.02, 1.72),
            (0.65, 3.67),
            (4.02, 3.67),
        ]
        for (_, row), (x, y) in zip(market.iterrows(), positions, strict=True):
            color = row["color"]
            add_rect(slide, x, y, 3.06, 1.62, fill="white", line="gray_200", radius=True)
            add_text(
                slide,
                row["name"],
                x + 0.20,
                y + 0.18,
                2.62,
                0.36,
                font_size=11,
                color="ink",
                bold=True,
            )
            add_text(
                slide,
                f"${row['tam'] / 1000:,.0f}B",
                x + 0.20,
                y + 0.64,
                1.35,
                0.50,
                font_size=24,
                color=color,
                bold=True,
            )
            add_tag(slide, "DERIVED", x + 2.06, y + 0.70, width=0.78)
            add_text(
                slide,
                f"{row['population_m']:,.0f}M units × ${row['annual_revenue_per_unit']:,.0f}/yr",
                x + 0.20,
                y + 1.23,
                2.62,
                0.22,
                font_size=7.5,
                color="gray_700",
            )

        add_rect(slide, 7.42, 1.72, 5.26, 3.57, fill="ink", radius=True)
        add_text(
            slide,
            "$250B",
            7.85,
            2.18,
            4.40,
            0.88,
            font_size=49,
            color="orange",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "total addressable market",
            7.85,
            3.05,
            4.40,
            0.34,
            font_size=14,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_line(slide, 8.30, 3.65, 11.80, 3.65, color="charcoal", width=1)
        add_text(
            slide,
            "$90B",
            8.05,
            3.98,
            1.70,
            0.55,
            font_size=28,
            color="blue",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "near-term SAM",
            9.75,
            4.08,
            1.90,
            0.28,
            font_size=10,
            color="gray_300",
            bold=True,
        )
        add_tag(slide, "EST.", 11.46, 4.08, width=0.62, fill="charcoal", color="white")

    def _slide_12_category_leadership(self) -> None:
        slide = self._new_slide(12)
        columns = [
            ("Replit", "orange"),
            ("AI IDEs", "blue"),
            ("App builders", "coral"),
            ("Model agents", "charcoal"),
            ("Cloud suites", "green"),
        ]
        rows = [
            ("Natural language → full app", [3, 2, 3, 2, 1]),
            ("Full code control", [3, 3, 2, 3, 3]),
            ("Real-time collaboration", [3, 2, 2, 1, 2]),
            ("Integrated cloud runtime", [3, 1, 2, 1, 3]),
            ("Enterprise governance", [3, 3, 1, 2, 3]),
            ("Consumer-scale distribution", [3, 1, 2, 2, 1]),
        ]
        table_x, table_y = 0.65, 1.72
        label_w, col_w = 3.22, 1.74
        header_h, row_h = 0.74, 0.60
        add_rect(slide, table_x, table_y, label_w + col_w * 5, header_h, fill="ink", radius=True)
        add_text(
            slide,
            "Observable platform capability",
            table_x + 0.22,
            table_y + 0.23,
            label_w - 0.35,
            0.24,
            font_size=9,
            color="white",
            bold=True,
        )
        for index, (label, color) in enumerate(columns):
            x = table_x + label_w + index * col_w
            if index == 0:
                add_rect(slide, x + 0.05, table_y + 0.08, col_w - 0.10, header_h - 0.16, fill="orange", radius=True)
            add_text(
                slide,
                label,
                x,
                table_y + 0.23,
                col_w,
                0.24,
                font_size=8,
                color="white",
                bold=True,
                align=PP_ALIGN.CENTER,
            )
        for row_index, (label, scores) in enumerate(rows):
            y = table_y + header_h + row_index * row_h
            fill = "white" if row_index % 2 == 0 else "gray_100"
            add_rect(slide, table_x, y, label_w + col_w * 5, row_h, fill=fill)
            add_text(
                slide,
                label,
                table_x + 0.22,
                y + 0.20,
                label_w - 0.35,
                0.24,
                font_size=8.5,
                color="ink",
                bold=True,
            )
            for col_index, score in enumerate(scores):
                cell_x = table_x + label_w + col_index * col_w
                add_matrix_dot(
                    slide,
                    cell_x + 0.63,
                    y + 0.26,
                    score,
                    active_color="orange" if col_index == 0 else columns[col_index][1],
                )
        add_rect(slide, 0.65, 6.12, 12.03, 0.40, fill="white", line="gray_200", radius=True)
        add_text(
            slide,
            "Replit covers more steps. One product brings in users, coordinates agents, and runs governed software.",
            0.92,
            6.20,
            11.50,
            0.20,
            font_size=8.5,
            color="gray_700",
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def _slide_13_historical_growth(self) -> None:
        slide = self._new_slide(13)
        financials = self.model.financials
        years = list(range(2023, 2029))
        revenue_scale = financials.loc[2028, "revenue"] / financials.loc[2024, "revenue"]
        add_stacked_columns(
            slide,
            0.65,
            2.10,
            7.60,
            4.23,
            [str(year) + ("E" if year >= 2026 else "A") for year in years],
            [
                {
                    "name": "Individual",
                    "values": financials.loc[years, "revenue_individual"].tolist(),
                    "color": "orange",
                },
                {
                    "name": "Enterprise",
                    "values": financials.loc[years, "revenue_enterprise"].tolist(),
                    "color": "blue",
                },
                {
                    "name": "Usage & deploy",
                    "values": financials.loc[years, "revenue_usage_and_deploy"].tolist(),
                    "color": "coral",
                },
            ],
            max_value=2250,
            value_format="${:,.0f}M",
        )
        add_text(
            slide,
            "REVENUE BY SOURCE, $M",
            0.80,
            1.72,
            3.0,
            0.25,
            font_size=8,
            color="gray_700",
            bold=True,
            font_name=FONT_MONO,
        )
        add_chart_legend(
            slide,
            [
                ("Individual", "orange"),
                ("Enterprise", "blue"),
                ("Agent use and hosting", "coral"),
            ],
            2.62,
            1.73,
            item_width=1.62,
        )
        add_tag(slide, "EST.", 7.35, 1.73, width=0.62)

        add_metric_card(
            slide,
            8.60,
            1.78,
            4.08,
            1.31,
            f"{revenue_scale:.0f}x",
            "2024 to 2028 revenue growth",
            accent="orange",
            evidence="EST.",
            value_size=28,
        )
        add_metric_card(
            slide,
            8.60,
            3.31,
            4.08,
            1.31,
            f"{financials.loc[2028, 'registered_users_m']:.0f}M",
            "Registered users by 2028E",
            accent="blue",
            evidence="EST.",
            value_size=28,
        )
        add_metric_card(
            slide,
            8.60,
            4.84,
            4.08,
            1.31,
            f"${financials.loc[2028, 'exit_arr'] / 1000:.2f}B",
            "2028E exit ARR",
            accent="coral",
            evidence="EST.",
            value_size=28,
        )
        add_text(
            slide,
            "Public anchor",
            0.90,
            6.31,
            0.95,
            0.20,
            font_size=7,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
        )
        add_text(
            slide,
            "Management stated Replit was on track for $1B run-rate revenue by year-end 2026.",
            1.88,
            6.29,
            6.40,
            0.22,
            font_size=8.2,
            color="gray_700",
        )

    def _slide_14_revenue_quality(self) -> None:
        slide = self._new_slide(14)
        rq = self.model.assumptions["revenue_quality"]
        add_text(
            slide,
            "REVENUE MIX",
            0.65,
            1.76,
            2.0,
            0.22,
            font_size=8,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
        )
        add_rect(slide, 0.65, 2.18, 7.25, 0.62, fill="gray_200", radius=True)
        add_rect(
            slide,
            0.65,
            2.18,
            7.25 * rq["recurring_share_2026"],
            0.62,
            fill="orange",
            radius=True,
        )
        add_rect(
            slide,
            0.65 + 7.25 * rq["recurring_share_2026"],
            2.18,
            7.25 * rq["usage_share_2026"],
            0.62,
            fill="blue",
            radius=True,
        )
        add_text(
            slide,
            f"{rq['recurring_share_2026']:.0%} recurring subscriptions",
            0.84,
            2.37,
            3.45,
            0.20,
            font_size=8,
            color="white",
            bold=True,
        )
        add_text(
            slide,
            f"{rq['usage_share_2026']:.0%} usage",
            5.94,
            2.37,
            1.60,
            0.20,
            font_size=8,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_tag(slide, "EST.", 7.98, 2.35, width=0.62)

        segment_values = [
            self.model.financials.loc[2026, "revenue_individual"],
            self.model.financials.loc[2026, "revenue_enterprise"],
            self.model.financials.loc[2026, "revenue_usage_and_deploy"],
        ]
        add_text(
            slide,
            "2026E REVENUE BY ENGINE",
            0.65,
            3.23,
            2.3,
            0.22,
            font_size=8,
            color="gray_700",
            bold=True,
            font_name=FONT_MONO,
        )
        add_horizontal_bars(
            slide,
            0.65,
            3.65,
            7.25,
            ["Individual", "Enterprise", "Usage & deploy"],
            segment_values,
            ["orange", "blue", "coral"],
            max_value=325,
            value_format="${:,.0f}M",
            row_height=0.70,
        )

        add_rect(slide, 8.43, 1.72, 4.25, 4.81, fill="white", line="gray_200", radius=True)
        add_text(
            slide,
            "GEOGRAPHIC MIX",
            8.74,
            2.02,
            2.0,
            0.22,
            font_size=8,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
        )
        geo = [
            ("United States", rq["us_share_2026"], "orange"),
            ("EMEA", rq["emea_share_2026"], "blue"),
            ("APAC", rq["apac_share_2026"], "coral"),
            ("Other", rq["other_share_2026"], "gray_500"),
        ]
        for index, (label, share, color) in enumerate(geo):
            y = 2.53 + index * 0.62
            add_rect(slide, 8.76, y, 2.55, 0.22, fill="gray_100", radius=True)
            add_rect(slide, 8.76, y, 2.55 * share / 0.60, 0.22, fill=color, radius=True)
            add_text(
                slide,
                label,
                8.76,
                y + 0.28,
                1.70,
                0.18,
                font_size=7,
                color="gray_700",
                bold=True,
            )
            add_text(
                slide,
                f"{share:.0%}",
                11.47,
                y - 0.01,
                0.58,
                0.24,
                font_size=8,
                color=color,
                bold=True,
                align=PP_ALIGN.RIGHT,
            )
        add_line(slide, 8.76, 5.17, 12.33, 5.17, color="gray_200", width=0.8)
        add_text(
            slide,
            f"{rq['top_10_customer_share_2026']:.0%}",
            8.76,
            5.48,
            0.82,
            0.38,
            font_size=21,
            color="orange",
            bold=True,
        )
        add_text(
            slide,
            "top-10 customer concentration",
            9.70,
            5.55,
            2.40,
            0.24,
            font_size=8.5,
            color="gray_700",
            bold=True,
        )
        add_tag(slide, "EST.", 11.73, 5.95, width=0.62)

    def _slide_15_retention(self) -> None:
        slide = self._new_slide(15)
        f = self.model.financials.loc[2026]
        add_text(
            slide,
            "ENTERPRISE ARR COHORT MULTIPLE",
            0.78,
            1.62,
            2.75,
            0.22,
            font_size=8,
            color="gray_700",
            bold=True,
            font_name=FONT_MONO,
        )
        add_native_line_chart(
            slide,
            0.65,
            1.95,
            7.55,
            4.23,
            ["Start", "M12", "M24", "M36"],
            [
                {
                    "name": "2023 cohort",
                    "values": [1.00, 1.12, 1.28, 1.42],
                    "color": "blue",
                    "end_label": "1.42x",
                },
                {
                    "name": "2024 cohort",
                    "values": [1.00, 1.22, 1.48, 1.70],
                    "color": "orange",
                    "end_label": "1.70x",
                },
            ],
            min_value=0.8,
            max_value=1.8,
            y_format="{:.1f}x",
        )
        add_tag(slide, "EST.", 7.42, 1.63, width=0.62)

        add_metric_card(
            slide,
            8.55,
            1.76,
            4.13,
            1.32,
            f"{f['enterprise_nrr']:.0%}",
            "2026E enterprise NRR",
            accent="orange",
            evidence="EST.",
            value_size=28,
        )
        add_metric_card(
            slide,
            8.55,
            3.31,
            4.13,
            1.32,
            f"{f['enterprise_grr']:.0%}",
            "2026E enterprise GRR",
            accent="blue",
            evidence="EST.",
            value_size=28,
        )
        add_rect(slide, 8.55, 4.86, 4.13, 1.30, fill="ink", radius=True)
        add_text(
            slide,
            "How accounts grow",
            8.88,
            5.10,
            2.10,
            0.26,
            font_size=11,
            color="white",
            bold=True,
        )
        add_bullet_list(
            slide,
            ["More builders", "More agents", "More deployed workloads"],
            8.88,
            5.47,
            3.25,
            font_size=7.8,
            color="gray_300",
            row_height=0.25,
        )
        add_rect(slide, 0.65, 6.28, 12.03, 0.24, fill="orange", radius=True)
        add_text(
            slide,
            "Retention is the main diligence gap. The cohort and retention data are estimates until the company provides actual data.",
            0.88,
            6.31,
            11.55,
            0.17,
            font_size=6.9,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def _slide_16_customer_proof(self) -> None:
        slide = self._new_slide(16)
        anchors = self.model.assumptions["public_anchors"]
        add_rect(slide, 0.65, 1.72, 5.02, 4.80, fill="ink", radius=True)
        add_text(
            slide,
            "UKG",
            0.98,
            2.03,
            1.4,
            0.42,
            font_size=23,
            color="white",
            bold=True,
        )
        add_tag(slide, "PUBLIC CASE STUDY", 3.80, 2.07, fill="orange", color="white", width=1.48)
        add_text(
            slide,
            "HR and workforce software | 16,000+ employees",
            1.02,
            2.47,
            3.90,
            0.24,
            font_size=8,
            color="gray_300",
        )
        add_text(
            slide,
            f"{anchors['ukg_feedback_capacity_multiple']:.0%}",
            0.98,
            2.82,
            2.20,
            0.70,
            font_size=42,
            color="orange",
            bold=True,
        )
        add_text(
            slide,
            "increase in feedback capacity",
            1.02,
            3.53,
            2.85,
            0.30,
            font_size=12,
            color="white",
            bold=True,
        )
        add_text(
            slide,
            "UKG built a reusable prototype system in Replit. Product and design teams now create working AI prototypes in days instead of weeks.",
            1.02,
            4.18,
            4.18,
            1.03,
            font_size=10,
            color="gray_300",
        )
        add_text(
            slide,
            "Rapid prototypes | Product work | Internal tools",
            1.02,
            5.62,
            3.95,
            0.30,
            font_size=8,
            color="coral",
            bold=True,
        )
        add_text(
            slide,
            "Source: Replit customer case study",
            1.02,
            6.02,
            3.98,
            0.24,
            font_size=7.8,
            color="gray_300",
        )

        add_rect(slide, 6.04, 1.72, 6.64, 2.22, fill="white", line="gray_200", radius=True)
        add_text(
            slide,
            "Enterprise references",
            6.36,
            2.02,
            2.4,
            0.26,
            font_size=12,
            color="ink",
            bold=True,
        )
        logos = [
            ("Atlassian", "orange"),
            ("PayPal", "blue"),
            ("Zillow", "coral"),
            ("Adobe", "charcoal"),
            ("Labcorp", "green"),
            ("Talkdesk", "orange"),
        ]
        for index, (name, color) in enumerate(logos):
            column = index % 3
            row = index // 3
            x = 6.36 + column * 2.00
            y = 2.52 + row * 0.60
            add_circle(slide, x, y + 0.02, 0.14, fill=color)
            add_text(
                slide,
                name,
                x + 0.25,
                y,
                1.35,
                0.24,
                font_size=8.5,
                color="ink",
                bold=True,
            )
        add_tag(slide, "NAMED BY REPLIT", 10.80, 3.49, width=1.45)

        add_rect(slide, 6.04, 4.27, 3.15, 2.25, fill="white", line="gray_200", radius=True)
        add_text(
            slide,
            "Minnesota Vikings",
            6.34,
            4.57,
            2.45,
            0.28,
            font_size=11,
            color="ink",
            bold=True,
        )
        add_text(
            slide,
            "Prototypes partnership ideas and saves staff time on game days.",
            6.34,
            5.09,
            2.45,
            0.68,
            font_size=9,
            color="gray_700",
        )
        add_tag(slide, "PUBLIC", 7.94, 6.05, width=0.92)

        add_rect(slide, 9.53, 4.27, 3.15, 2.25, fill="white", line="gray_200", radius=True)
        add_text(
            slide,
            "Zillow",
            9.83,
            4.57,
            2.45,
            0.28,
            font_size=11,
            color="ink",
            bold=True,
        )
        add_text(
            slide,
            "Uses live collaboration with partners to turn immediate feedback into measurable product wins.",
            9.83,
            5.09,
            2.45,
            0.68,
            font_size=9,
            color="gray_700",
        )
        add_tag(slide, "PUBLIC", 11.43, 6.05, width=0.92)

    def _slide_17_business_model(self) -> None:
        slide = self._new_slide(17)
        stages = [
            ("FREE", "Acquire", "50M+ users", "orange", 7.00),
            ("CORE", "Convert", "plans + credits", "blue", 6.30),
            ("TEAMS", "Expand", "collaboration", "coral", 5.60),
            ("ENTERPRISE", "Commit", "committed usage", "green", 4.90),
            ("RUNTIME", "Grow", "deploy + operate", "charcoal", 4.20),
        ]
        x = 0.65
        for index, (stage, action, detail, color, width) in enumerate(stages):
            y = 1.95 + index * 0.86
            bar_x = x + (7.00 - width) / 2
            add_rect(slide, bar_x, y, width, 0.64, fill=color, radius=True)
            add_text(
                slide,
                stage,
                bar_x + 0.20,
                y + 0.18,
                1.05,
                0.22,
                font_size=8,
                color="white",
                bold=True,
                font_name=FONT_MONO,
            )
            add_text(
                slide,
                action,
                bar_x + 1.30,
                y + 0.18,
                1.08,
                0.22,
                font_size=9,
                color="white",
                bold=True,
            )
            add_text(
                slide,
                detail,
                bar_x + 2.42,
                y + 0.18,
                width - 2.64,
                0.22,
                font_size=7.3,
                color="white",
                align=PP_ALIGN.RIGHT,
            )

        add_rect(slide, 8.34, 1.72, 4.34, 4.80, fill="white", line="gray_200", radius=True)
        add_text(
            slide,
            "REVENUE SOURCES",
            8.69,
            2.02,
            3.65,
            0.24,
            font_size=8,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
            align=PP_ALIGN.CENTER,
        )
        vectors = [
            ("Subscription", "Access, features, collaboration", "orange"),
            ("Agent usage", "Outcome-linked credits", "blue"),
            ("Enterprise", "Committed spend + governance", "coral"),
            ("Runtime", "Compute, database, deployment", "green"),
        ]
        for index, (title, body, color) in enumerate(vectors):
            y = 2.52 + index * 0.83
            add_circle(slide, 8.73, y, 0.42, fill=color)
            add_text(
                slide,
                str(index + 1),
                8.73,
                y,
                0.42,
                0.42,
                font_size=8,
                color="white",
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                font_name=FONT_MONO,
            )
            add_text(
                slide,
                title,
                9.35,
                y - 0.01,
                1.35,
                0.22,
                font_size=9,
                color="ink",
                bold=True,
            )
            add_text(
                slide,
                body,
                10.67,
                y - 0.01,
                1.58,
                0.36,
                font_size=7.5,
                color="gray_700",
            )
        add_rect(slide, 8.70, 5.89, 3.62, 0.36, fill="ink", radius=True)
        add_text(
            slide,
            "Price scales with customer success",
            8.70,
            5.97,
            3.62,
            0.20,
            font_size=8,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def _slide_18_unit_economics(self) -> None:
        slide = self._new_slide(18)
        unit = self.model.unit_economics
        years = list(unit.index)
        add_text(
            slide,
            "GROSS-MARGIN EXPANSION",
            0.78,
            1.62,
            2.45,
            0.22,
            font_size=8,
            color="gray_700",
            bold=True,
            font_name=FONT_MONO,
        )
        add_native_line_chart(
            slide,
            0.65,
            1.95,
            5.60,
            3.18,
            [str(year) + "E" for year in years],
            [
                {
                    "name": "Gross margin",
                    "values": (unit["blended_gross_margin"] * 100).tolist(),
                    "color": "orange",
                    "end_label": "71%",
                }
            ],
            min_value=20,
            max_value=80,
            y_format="{:.0f}%",
            show_legend=False,
        )
        add_text(
            slide,
            "COMPUTE COST / SUCCESSFUL BUILD",
            6.83,
            1.62,
            3.10,
            0.22,
            font_size=8,
            color="gray_700",
            bold=True,
            font_name=FONT_MONO,
        )
        add_native_line_chart(
            slide,
            6.70,
            1.95,
            5.98,
            3.18,
            [str(year) + "E" for year in years],
            [
                {
                    "name": "Compute index",
                    "values": unit["compute_cost_per_successful_build_index"].tolist(),
                    "color": "blue",
                    "end_label": "35",
                }
            ],
            min_value=20,
            max_value=110,
            y_format="{:.0f}",
            show_legend=False,
        )
        metrics = [
            (f"{unit.loc[2026, 'individual_payback_months']:.0f} mo.", "Individual CAC payback", "orange"),
            (f"{unit.loc[2026, 'enterprise_payback_months']:.0f} mo.", "Enterprise CAC payback", "blue"),
            (f"{unit.loc[2026, 'enterprise_ltv_cac']:.1f}x", "Enterprise LTV / CAC", "coral"),
        ]
        for index, (value, label, color) in enumerate(metrics):
            add_metric_card(
                slide,
                0.65 + index * 4.10,
                5.28,
                3.77,
                1.21,
                value,
                label,
                accent=color,
                evidence="EST.",
                value_size=23,
            )

    def _slide_19_gtm(self) -> None:
        slide = self._new_slide(19)
        stages = [
            ("CREATE", "Individual builder", "Free + paid", "orange"),
            ("SHARE", "Team adoption", "Collaboration", "blue"),
            ("DEPLOY", "Business workload", "Usage", "coral"),
            ("GOVERN", "Enterprise standard", "Commitment", "green"),
            ("EXPAND", "More teams + apps", "NRR", "charcoal"),
        ]
        for index, (stage, buyer, motion, color) in enumerate(stages):
            x = 0.65 + index * 2.44
            add_circle(slide, x + 0.57, 2.12, 1.26, fill=color)
            add_text(
                slide,
                stage,
                x + 0.57,
                2.12,
                1.26,
                1.26,
                font_size=9,
                color="white",
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                font_name=FONT_MONO,
            )
            if index < 4:
                add_chevron(slide, x + 1.94, 2.55, 0.30, 0.40, fill="gray_300")
            add_text(
                slide,
                buyer,
                x,
                3.67,
                2.40,
                0.27,
                font_size=10,
                color="ink",
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            add_tag(
                slide,
                motion,
                x + 0.60,
                4.12,
                fill="white",
                color=color,
                width=1.20,
            )

        add_rect(slide, 0.65, 4.86, 12.03, 1.63, fill="ink", radius=True)
        add_text(
            slide,
            "Product-led growth lowers acquisition cost",
            0.98,
            5.17,
            3.42,
            0.48,
            font_size=12,
            color="white",
            bold=True,
        )
        add_text(
            slide,
            "Enterprise sales monetizes organic demand",
            4.88,
            5.17,
            3.42,
            0.48,
            font_size=12,
            color="white",
            bold=True,
        )
        add_text(
            slide,
            "Runtime increases account spend",
            8.83,
            5.17,
            3.10,
            0.48,
            font_size=12,
            color="white",
            bold=True,
        )
        for x in (4.57, 8.53):
            add_line(slide, x, 5.10, x, 6.17, color="charcoal", width=1)
        add_text(
            slide,
            "50M+ user base",
            0.98,
            5.82,
            3.42,
            0.24,
            font_size=9,
            color="orange",
            bold=True,
        )
        add_text(
            slide,
            "Usage signals identify the best accounts",
            4.88,
            5.82,
            3.42,
            0.24,
            font_size=9,
            color="blue",
            bold=True,
        )
        add_text(
            slide,
            "Apps + agents + compute",
            8.83,
            5.82,
            3.10,
            0.24,
            font_size=9,
            color="coral",
            bold=True,
        )

    def _slide_20_pipeline(self) -> None:
        slide = self._new_slide(20)
        pipeline = self.model.assumptions["pipeline"]
        funnel = [
            ("Qualified pipeline", pipeline["qualified_pipeline"], "orange"),
            ("Late-stage pipeline", pipeline["late_stage_pipeline"], "blue"),
            ("Weighted pipeline", pipeline["weighted_pipeline"], "coral"),
            ("New ARR target", pipeline["next_12m_new_arr_target"], "green"),
        ]
        center = 4.05
        for index, (label, value, color) in enumerate(funnel):
            width = 6.75 * value / funnel[0][1]
            x = center - width / 2 + 0.65
            y = 1.82 + index * 1.03
            add_rect(slide, x, y, width, 0.72, fill=color, radius=True)
            add_text(
                slide,
                label,
                x + 0.16,
                y + 0.21,
                max(1.25, width - 1.30),
                0.24,
                font_size=8,
                color="white",
                bold=True,
            )
            add_text(
                slide,
                f"${value:,.0f}M",
                x + width - 1.12,
                y + 0.19,
                0.92,
                0.28,
                font_size=10,
                color="white",
                bold=True,
                align=PP_ALIGN.RIGHT,
            )
        add_tag(slide, "EST.", 7.48, 1.86, width=0.62)
        add_text(
            slide,
            "ENTERPRISE FUNNEL",
            0.65,
            1.76,
            2.0,
            0.22,
            font_size=8,
            color="gray_700",
            bold=True,
            font_name=FONT_MONO,
        )

        cards = [
            (f"{pipeline['coverage']:.1f}x", "Pipeline coverage", "orange"),
            (f"{pipeline['win_rate']:.0%}", "Win rate", "blue"),
            (f"{pipeline['median_sales_cycle_days']}d", "Median sales cycle", "coral"),
            (f"{pipeline['quota_attainment']:.0%}", "Quota attainment", "green"),
        ]
        for index, (value, label, color) in enumerate(cards):
            add_metric_card(
                slide,
                8.38,
                1.72 + index * 1.17,
                4.30,
                0.98,
                value,
                label,
                accent=color,
                evidence="EST.",
                value_size=19,
            )
        add_rect(slide, 0.65, 6.08, 12.03, 0.44, fill="white", line="gray_200", radius=True)
        add_text(
            slide,
            "Execution focus: improve late-stage conversion, shorten security review, and lift productivity as the enterprise sales force matures.",
            0.92,
            6.18,
            11.50,
            0.23,
            font_size=8.8,
            color="gray_700",
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def _slide_21_competition(self) -> None:
        slide = self._new_slide(21)
        add_line(slide, 1.58, 5.98, 11.95, 5.98, color="gray_500", width=1.1)
        add_line(slide, 1.58, 1.88, 1.58, 5.98, color="gray_500", width=1.1)
        add_text(
            slide,
            "Point tool",
            0.76,
            5.82,
            0.70,
            0.22,
            font_size=7,
            color="gray_700",
            align=PP_ALIGN.RIGHT,
        )
        add_text(
            slide,
            "Full platform",
            0.70,
            1.76,
            0.76,
            0.22,
            font_size=7,
            color="gray_700",
            align=PP_ALIGN.RIGHT,
        )
        add_text(
            slide,
            "Professional developers",
            1.58,
            6.14,
            1.95,
            0.22,
            font_size=7,
            color="gray_700",
        )
        add_text(
            slide,
            "Everyone",
            10.98,
            6.14,
            0.95,
            0.22,
            font_size=7,
            color="gray_700",
            align=PP_ALIGN.RIGHT,
        )
        add_text(
            slide,
            "PLATFORM BREADTH",
            0.18,
            3.63,
            1.18,
            0.22,
            font_size=7,
            color="gray_700",
            bold=True,
            font_name=FONT_MONO,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "AUDIENCE BREADTH",
            5.60,
            6.43,
            2.15,
            0.22,
            font_size=7,
            color="gray_700",
            bold=True,
            font_name=FONT_MONO,
            align=PP_ALIGN.CENTER,
        )
        competitors = [
            ("GitHub\nCopilot", 3.25, 4.86, "charcoal"),
            ("Cursor /\nWindsurf", 3.15, 3.66, "blue"),
            ("Claude\nCode", 2.36, 4.08, "coral"),
            ("Lovable /\nBolt / v0", 8.58, 4.26, "coral"),
            ("Cloud\nsuites", 3.55, 2.28, "green"),
            ("Replit", 9.48, 2.04, "orange"),
        ]
        for name, x, y, color in competitors:
            diameter = 1.15 if name == "Replit" else 0.88
            add_circle(slide, x, y, diameter, fill=color)
            add_text(
                slide,
                name,
                x,
                y,
                diameter,
                diameter,
                font_size=9 if name == "Replit" else 7.2,
                color="white",
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
            )
        add_rect(slide, 7.74, 1.72, 4.94, 0.46, fill="orange", radius=True)
        add_text(
            slide,
            "Replit serves more users and covers the full workflow",
            7.90,
            1.83,
            4.62,
            0.22,
            font_size=8.5,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_tag(slide, "OBSERVABLE POSITIONING", 10.59, 5.45, width=1.68)

    def _slide_22_moat(self) -> None:
        slide = self._new_slide(22)
        center_x, center_y = 5.55, 2.24
        add_circle(slide, center_x, center_y, 2.26, fill="ink")
        add_text(
            slide,
            "MORE\nSUCCESSFUL\nBUILDS",
            center_x,
            center_y,
            2.26,
            2.26,
            font_size=15,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            font_name=FONT_MONO,
        )
        moat = [
            ("Distribution", "50M+ users lower acquisition cost", "orange", (0.65, 1.72)),
            ("Workflow data", "Plans, edits, tests, and repairs improve agents", "blue", (8.86, 1.72)),
            ("Runtime", "Deployed apps add revenue and product data", "coral", (0.65, 4.62)),
            ("Integrations", "Enterprise systems make Replit more useful", "green", (8.86, 4.62)),
            ("Governance", "Identity and policy increase switching cost", "charcoal", (4.70, 5.32)),
        ]
        for index, (title, body, color, (x, y)) in enumerate(moat, 1):
            add_rect(slide, x, y, 3.82, 1.31, fill="white", line="gray_200", radius=True)
            add_circle(slide, x + 0.22, y + 0.24, 0.46, fill=color)
            add_text(
                slide,
                str(index),
                x + 0.22,
                y + 0.24,
                0.46,
                0.46,
                font_size=8,
                color="white",
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                font_name=FONT_MONO,
            )
            add_text(
                slide,
                title,
                x + 0.84,
                y + 0.20,
                2.55,
                0.28,
                font_size=11,
                color="ink",
                bold=True,
            )
            add_text(
                slide,
                body,
                x + 0.84,
                y + 0.58,
                2.55,
                0.44,
                font_size=8.2,
                color="gray_700",
            )
        connectors = [
            (4.47, 2.38, 5.55, 3.17),
            (8.86, 2.38, 7.81, 3.17),
            (4.47, 5.27, 5.68, 4.44),
            (8.86, 5.27, 7.66, 4.44),
            (6.61, 5.32, 6.67, 4.50),
        ]
        for x1, y1, x2, y2 in connectors:
            add_line(slide, x1, y1, x2, y2, color="gray_300", width=1.2)

    def _slide_23_historical_financials(self) -> None:
        slide = self._new_slide(23)
        financials = self.model.financials
        years = list(range(2025, 2032))
        metrics = [
            ("Revenue", "revenue", "money"),
            ("Growth", "revenue_growth", "percent"),
            ("Gross profit", "gross_profit", "money"),
            ("Gross margin", "gross_margin", "percent"),
            ("R&D", "r_and_d", "money"),
            ("Sales and marketing", "sales_and_marketing", "money"),
            ("G&A", "g_and_a", "money"),
            ("Adjusted operating income", "adjusted_operating_income", "money"),
            ("Adjusted operating margin", "adjusted_operating_margin", "percent"),
            ("Free cash flow", "free_cash_flow", "money"),
        ]
        table_x, table_y = 0.65, 1.72
        label_w = 2.40
        col_w = (12.03 - label_w) / len(years)
        header_h = 0.52
        row_h = 0.38
        add_rect(
            slide,
            table_x,
            table_y,
            label_w + col_w * len(years),
            header_h,
            fill="ink",
            radius=True,
        )
        add_text(
            slide,
            "P&L MODEL",
            table_x + 0.22,
            table_y + 0.16,
            label_w - 0.35,
            0.24,
            font_size=8,
            color="white",
            bold=True,
            font_name=FONT_MONO,
        )
        for index, year in enumerate(years):
            add_text(
                slide,
                f"{year}E",
                table_x + label_w + index * col_w,
                table_y + 0.16,
                col_w,
                0.24,
                font_size=8,
                color="white",
                bold=True,
                align=PP_ALIGN.CENTER,
            )
        for row_index, (label, column, fmt) in enumerate(metrics):
            y = table_y + header_h + row_index * row_h
            fill = "white" if row_index % 2 == 0 else "gray_100"
            add_rect(
                slide,
                table_x,
                y,
                label_w + col_w * len(years),
                row_h,
                fill=fill,
            )
            add_text(
                slide,
                label,
                table_x + 0.22,
                y + 0.09,
                label_w - 0.35,
                0.24,
                font_size=8,
                color="ink",
                bold=row_index in (0, 2, 7, 9),
            )
            for year_index, year in enumerate(years):
                value = financials.loc[year, column]
                if fmt == "percent":
                    display = "n/a" if value != value else f"{value:.0%}"
                else:
                    display = (
                        f"(${abs(value):,.0f})" if value < 0 else f"${value:,.0f}"
                    )
                value_color = (
                    "orange"
                    if column in ("adjusted_operating_income", "free_cash_flow")
                    and value > 0
                    else "ink"
                )
                add_text(
                    slide,
                    display,
                    table_x + label_w + year_index * col_w,
                    y + 0.09,
                    col_w,
                    0.24,
                    font_size=7.5,
                    color=value_color,
                    bold=row_index in (0, 2, 7, 9),
                    align=PP_ALIGN.CENTER,
                )
        add_tag(slide, "EST.", 11.85, 1.82, width=0.62)
        add_rect(slide, 0.65, 6.16, 12.03, 0.36, fill="orange", radius=True)
        add_text(
            slide,
            "Adjusted operating income and free cash flow turn positive in 2028.",
            0.90,
            6.23,
            11.55,
            0.20,
            font_size=8.5,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def _slide_24_operating_plan(self) -> None:
        slide = self._new_slide(24)
        financials = self.model.financials
        years = list(range(2026, 2032))
        revenue_cagr = (
            financials.loc[2031, "revenue"] / financials.loc[2026, "revenue"]
        ) ** (1 / 5) - 1
        enterprise_runtime_mix = (
            financials.loc[2031, "revenue_enterprise"]
            + financials.loc[2031, "revenue_usage_and_deploy"]
        ) / financials.loc[2031, "revenue"]
        add_stacked_columns(
            slide,
            0.65,
            2.10,
            7.60,
            4.23,
            [str(year) + "E" for year in years],
            [
                {
                    "name": "Individual",
                    "values": financials.loc[years, "revenue_individual"].tolist(),
                    "color": "orange",
                },
                {
                    "name": "Enterprise",
                    "values": financials.loc[years, "revenue_enterprise"].tolist(),
                    "color": "blue",
                },
                {
                    "name": "Usage & deploy",
                    "values": financials.loc[years, "revenue_usage_and_deploy"].tolist(),
                    "color": "coral",
                },
            ],
            max_value=5800,
            value_format="${:,.0f}M",
        )
        add_text(
            slide,
            "REVENUE BY SOURCE, $M",
            0.80,
            1.72,
            2.25,
            0.22,
            font_size=8,
            color="gray_700",
            bold=True,
            font_name=FONT_MONO,
        )
        add_chart_legend(
            slide,
            [
                ("Individual", "orange"),
                ("Enterprise", "blue"),
                ("Agent use and hosting", "coral"),
            ],
            2.62,
            1.73,
            item_width=1.62,
        )
        add_tag(slide, "EST.", 7.35, 1.73, width=0.62)

        metrics = [
            (f"${financials.loc[2031, 'revenue'] / 1000:.2f}B", "2031E revenue", "orange"),
            (f"{revenue_cagr:.0%}", "2026 to 2031 CAGR", "blue"),
            (f"{financials.loc[2031, 'gross_margin']:.0%}", "2031E gross margin", "coral"),
            (f"{financials.loc[2031, 'international_revenue_share']:.0%}", "2031E international mix", "green"),
        ]
        for index, (value, label, color) in enumerate(metrics):
            add_metric_card(
                slide,
                8.60,
                1.78 + index * 1.12,
                4.08,
                0.93,
                value,
                label,
                accent=color,
                evidence="EST.",
                value_size=18,
            )
        add_rect(slide, 8.60, 6.31, 4.08, 0.21, fill="orange", radius=True)
        add_text(
            slide,
            f"Enterprise + runtime = {enterprise_runtime_mix:.0%} of 2031 revenue",
            8.60,
            6.29,
            4.08,
            0.19,
            font_size=7.2,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def _slide_25_profitability(self) -> None:
        slide = self._new_slide(25)
        financials = self.model.financials
        years = list(range(2026, 2032))
        add_native_line_chart(
            slide,
            0.65,
            1.78,
            7.35,
            4.60,
            [str(year) + "E" for year in years],
            [
                {
                    "name": "Gross margin",
                    "values": (financials.loc[years, "gross_margin"] * 100).tolist(),
                    "color": "orange",
                    "end_label": "76%",
                },
                {
                    "name": "Adj. op. margin",
                    "values": (
                        financials.loc[years, "adjusted_operating_margin"] * 100
                    ).tolist(),
                    "color": "blue",
                    "end_label": "31%",
                },
                {
                    "name": "FCF margin",
                    "values": (
                        financials.loc[years, "free_cash_flow_margin"] * 100
                    ).tolist(),
                    "color": "coral",
                    "end_label": "31%",
                },
            ],
            min_value=-80,
            max_value=80,
            y_format="{:.0f}%",
        )
        add_tag(slide, "EST.", 7.16, 1.79, width=0.62)
        add_rect(slide, 8.34, 1.78, 4.34, 4.60, fill="white", line="gray_200", radius=True)
        add_text(
            slide,
            "MARGIN DRIVERS",
            8.68,
            2.08,
            3.66,
            0.24,
            font_size=8,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
            align=PP_ALIGN.CENTER,
        )
        drivers = [
            ("+30 pts", "Gross-margin expansion", "orange"),
            ("+11 pts", "Sales productivity", "blue"),
            ("+9 pts", "G&A efficiency", "coral"),
            ("(4) pts", "Ongoing R&D reinvestment", "green"),
        ]
        for index, (value, label, color) in enumerate(drivers):
            y = 2.60 + index * 0.76
            add_text(
                slide,
                value,
                8.68,
                y,
                0.92,
                0.28,
                font_size=13,
                color=color,
                bold=True,
            )
            add_text(
                slide,
                label,
                9.77,
                y + 0.03,
                2.27,
                0.24,
                font_size=8.5,
                color="ink",
                bold=True,
            )
        add_line(slide, 8.68, 5.72, 12.30, 5.72, color="gray_200", width=0.8)
        add_text(
            slide,
            "2028E",
            8.68,
            5.92,
            0.84,
            0.26,
            font_size=8,
            color="gray_700",
            font_name=FONT_MONO,
        )
        add_text(
            slide,
            "first full year of positive free cash flow",
            9.56,
            5.90,
            2.65,
            0.28,
            font_size=9,
            color="ink",
            bold=True,
        )

    def _slide_26_use_of_proceeds(self) -> None:
        slide = self._new_slide(26)
        uses = self.model.use_of_proceeds
        tx = self.model.transaction
        short_labels = {
            "Model training and infrastructure": "Model training",
            "Enterprise expansion": "Enterprise sales",
            "User growth": "User growth",
            "International expansion": "International",
            "Talent and research": "Talent and research",
            "Execution reserve": "Reserve",
        }
        add_text(
            slide,
            "PRIMARY CAPITAL",
            0.65,
            1.77,
            1.75,
            0.22,
            font_size=8,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
        )
        add_horizontal_bars(
            slide,
            0.65,
            2.15,
            7.00,
            [short_labels[value] for value in uses["category"]],
            uses["amount"].tolist(),
            uses["color"].tolist(),
            max_value=230,
            value_format="${:,.0f}M",
            row_height=0.65,
        )
        add_rect(slide, 8.08, 1.72, 4.60, 4.82, fill="ink", radius=True)
        add_text(
            slide,
            f"${tx['primary']:.0f}M",
            8.48,
            2.18,
            3.80,
            0.65,
            font_size=39,
            color="orange",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "company proceeds",
            8.48,
            2.86,
            3.80,
            0.30,
            font_size=12,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "WHAT THE CAPITAL FUNDS",
            8.48,
            3.40,
            3.80,
            0.24,
            font_size=8,
            color="orange",
            bold=True,
            align=PP_ALIGN.CENTER,
            font_name=FONT_MONO,
        )
        add_bullet_list(
            slide,
            [
                "65% lower compute cost per successful build",
                "4,200 enterprise customers",
                "Local infrastructure, sales, and compliance in Europe, India, and the Middle East",
                "100M users and 5M paid users",
                "Model, security, product, and sales leaders",
            ],
            8.62,
            3.82,
            3.50,
            font_size=8,
            color="gray_300",
            row_height=0.38,
        )
        add_line(slide, 8.67, 5.78, 12.09, 5.78, color="charcoal", width=1)
        add_text(
            slide,
            f"${tx['secondary']:.0f}M secondary",
            8.48,
            5.92,
            2.10,
            0.26,
            font_size=11,
            color="coral",
            bold=True,
            align=PP_ALIGN.LEFT,
        )
        add_text(
            slide,
            "No cash to the company",
            10.45,
            5.93,
            1.48,
            0.24,
            font_size=8,
            color="gray_300",
        )

    def _slide_27_milestones(self) -> None:
        slide = self._new_slide(27)
        cash = self.model.monthly_cash_flow.loc[
            : self.model.cash_summary["ipo_target_date"]
        ]
        summary = self.model.cash_summary
        labels = [date.strftime("%b-%y") for date in cash.index]
        no_series_e = [
            value if date <= summary["funding_need_date"] else None
            for date, value in cash["ending_cash_without_series_e"].items()
        ]
        add_text(
            slide,
            "ENDING CASH BY MONTH, $M",
            0.78,
            1.72,
            2.65,
            0.22,
            font_size=8,
            color="gray_700",
            bold=True,
            font_name=FONT_MONO,
        )
        add_native_line_chart(
            slide,
            0.65,
            2.04,
            8.10,
            4.05,
            labels,
            [
                {
                    "name": "Base with Series E",
                    "values": cash["ending_cash"].tolist(),
                    "color": "orange",
                    "end_label": f"${cash['ending_cash'].iloc[-1]:,.0f}M",
                },
                {
                    "name": "Downside with Series E",
                    "values": cash["ending_cash_downside"].tolist(),
                    "color": "coral",
                    "end_label": f"${cash['ending_cash_downside'].iloc[-1]:,.0f}M",
                },
                {
                    "name": "No Series E",
                    "values": no_series_e,
                    "color": "gray_500",
                    "end_label": f"Below ${summary['minimum_cash']:.0f}M",
                },
                {
                    "name": "Minimum cash",
                    "values": cash["minimum_cash"].tolist(),
                    "color": "orange_dark",
                    "end_label": f"${summary['minimum_cash']:.0f}M",
                    "dash": True,
                    "width": 1.2,
                },
            ],
            min_value=0,
            max_value=3000,
            y_format="${:,.0f}",
            x_label_every=6,
            show_markers=False,
        )
        plot_left = 1.13
        plot_width = 7.38
        funding_index = cash.index.get_loc(summary["financing_close_date"])
        funding_x = plot_left + funding_index * plot_width / (len(cash) - 1)
        add_line(
            slide,
            funding_x,
            2.39,
            funding_x,
            5.69,
            color="orange_dark",
            width=1.1,
            dash=True,
        )
        add_text(
            slide,
            f"{summary['financing_close_date'].strftime('%b-%y')} funding",
            funding_x + 0.08,
            6.05,
            1.05,
            0.20,
            font_size=6.5,
            color="orange_dark",
            bold=True,
        )
        ipo_x = plot_left + plot_width
        add_line(
            slide,
            ipo_x,
            2.39,
            ipo_x,
            5.69,
            color="orange",
            width=1.1,
            dash=True,
        )
        add_text(
            slide,
            "H2-30 IPO",
            ipo_x - 0.78,
            6.05,
            0.75,
            0.20,
            font_size=6.5,
            color="orange",
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
        add_tag(slide, "MONTHLY MODEL", 7.43, 1.72, width=1.18)

        cards = [
            (
                f"${summary['beginning_cash']:.0f}M",
                "Cash at Sep-26",
                "orange",
            ),
            (
                summary["funding_need_date"].strftime("%b-%y"),
                f"Below ${summary['minimum_cash']:.0f}M without the round",
                "blue",
            ),
            (
                "No round",
                "Further equity before IPO in base case",
                "coral",
            ),
            (
                summary["ipo_target_date"].strftime("H2 %Y"),
                "Target IPO window",
                "green",
            ),
        ]
        for index, (value, label, color) in enumerate(cards):
            add_metric_card(
                slide,
                9.10,
                1.72 + index * 1.13,
                3.58,
                0.95,
                value,
                label,
                accent=color,
                evidence="EST.",
                value_size=18,
            )
        add_rect(slide, 9.10, 6.27, 3.58, 0.25, fill="orange", radius=True)
        add_text(
            slide,
            "The base and downside cases fund the company through the IPO target.",
            9.10,
            6.30,
            3.58,
            0.18,
            font_size=7,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def _slide_28_leadership(self) -> None:
        slide = self._new_slide(28)
        add_rect(slide, 0.65, 1.72, 3.66, 4.80, fill="ink", radius=True)
        add_circle(slide, 1.04, 2.14, 1.14, fill="orange")
        add_text(
            slide,
            "AM",
            1.04,
            2.14,
            1.14,
            1.14,
            font_size=24,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            font_name=FONT_MONO,
        )
        add_text(
            slide,
            "Amjad Masad",
            1.04,
            3.59,
            2.80,
            0.38,
            font_size=18,
            color="white",
            bold=True,
        )
        add_text(
            slide,
            "Founder & CEO",
            1.04,
            4.05,
            2.30,
            0.26,
            font_size=9,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
        )
        add_text(
            slide,
            "Founder-led product conviction anchors the long-term mission: make software creation accessible to everyone.",
            1.04,
            4.68,
            2.72,
            0.92,
            font_size=9.5,
            color="gray_300",
        )
        add_tag(slide, "PUBLIC", 2.95, 5.91, fill="charcoal", color="white", width=0.92)

        add_text(
            slide,
            "SERIES E OPERATING MANDATE",
            4.76,
            1.76,
            3.2,
            0.24,
            font_size=8,
            color="orange",
            bold=True,
            font_name=FONT_MONO,
        )
        priorities = [
            ("Financial rigor", "Segment ownership, forecasting cadence, and unit-economics accountability.", "orange"),
            ("Enterprise execution", "Scaled field leadership, partner motion, and customer-success depth.", "blue"),
            ("Global operations", "Regional product, infrastructure, compliance, and market leadership.", "coral"),
            ("Governance", "Add independent directors, audit systems, risk controls, security controls, and compensation processes.", "green"),
        ]
        for index, (title, body, color) in enumerate(priorities):
            row = index % 2
            column = index // 2
            x = 4.70 + column * 4.03
            y = 2.21 + row * 1.88
            add_rect(slide, x, y, 3.64, 1.55, fill="white", line="gray_200", radius=True)
            add_text(
                slide,
                f"0{index + 1}",
                x + 0.20,
                y + 0.18,
                0.34,
                0.22,
                font_size=8,
                color=color,
                bold=True,
                font_name=FONT_MONO,
            )
            add_text(
                slide,
                title,
                x + 0.20,
                y + 0.52,
                3.20,
                0.30,
                font_size=11,
                color="ink",
                bold=True,
            )
            add_text(
                slide,
                body,
                x + 0.20,
                y + 0.92,
                3.20,
                0.42,
                font_size=8,
                color="gray_700",
            )
        add_rect(slide, 4.70, 6.14, 7.98, 0.38, fill="gray_100", radius=True)
        add_text(
            slide,
            "The round underwrites management-system scale without diluting founder-led product velocity.",
            4.92,
            6.22,
            7.54,
            0.22,
            font_size=8.2,
            color="gray_700",
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def _slide_29_risks(self) -> None:
        slide = self._new_slide(29)
        rows = [
            ("Model dependency and cost", "High", "Routing, Replit model training, capacity contracts, and usage pricing"),
            ("Competition", "High", "Own the workflow, runtime, collaboration, and distribution"),
            ("Code quality and security", "High", "Automated testing, Security Agent, policy, audit logs, and human review"),
            ("Enterprise conversion", "Medium", "Product-led sales, field sales, partners, customer success, and governance"),
            ("Revenue concentration", "Medium", "A broad user base, low estimated account concentration, and global sales"),
            ("International delivery", "Medium", "Regional infrastructure, compliance, payments, and channel partners"),
        ]
        table_x, table_y = 0.65, 1.72
        widths = [2.75, 1.20, 8.08]
        headers = ["Risk", "Impact", "Mitigation"]
        cursor = table_x
        add_rect(slide, table_x, table_y, sum(widths), 0.68, fill="ink", radius=True)
        for width, header in zip(widths, headers, strict=True):
            add_text(
                slide,
                header,
                cursor + 0.16,
                table_y + 0.21,
                width - 0.28,
                0.24,
                font_size=8,
                color="white",
                bold=True,
                font_name=FONT_MONO,
                align=PP_ALIGN.CENTER if header == "Impact" else PP_ALIGN.LEFT,
            )
            cursor += width
        severity_color = {"High": "red", "Medium": "yellow"}
        for row_index, row in enumerate(rows):
            y = table_y + 0.68 + row_index * 0.66
            fill = "white" if row_index % 2 == 0 else "gray_100"
            add_rect(slide, table_x, y, sum(widths), 0.66, fill=fill)
            cursor = table_x
            for col_index, (value, width) in enumerate(zip(row, widths, strict=True)):
                if col_index == 1:
                    color = severity_color[value]
                    add_tag(
                        slide,
                        value,
                        cursor + (width - 0.72) / 2,
                        y + 0.20,
                        fill=color,
                        color="white",
                        width=0.72,
                    )
                else:
                    add_text(
                        slide,
                        value,
                        cursor + 0.16,
                        y + 0.18,
                        width - 0.28,
                        0.30,
                        font_size=8.2 if col_index == 0 else 7.6,
                        color="ink" if col_index == 0 else "gray_700",
                        bold=col_index == 0,
                    )
                cursor += width

    def _slide_30_closing(self) -> None:
        slide = self._new_slide(30, dark=True, show_header=False)
        financials = self.model.financials
        total_tam = self.model.market["tam"].sum()
        add_replit_mark(
            slide,
            0.76,
            0.62,
            scale=0.74,
            wordmark=True,
            color="orange",
            wordmark_color="white",
        )
        add_tag(
            slide,
            "Series E",
            11.46,
            0.68,
            fill="orange",
            color="white",
            width=1.05,
        )
        add_text(
            slide,
            "Replit can become the default way\nthe world creates software",
            0.76,
            1.62,
            10.45,
            1.60,
            font_size=34,
            color="white",
            bold=True,
            line_spacing=0.91,
        )
        thesis = [
            (
                "01",
                "Mass distribution",
                f"{financials.loc[2026, 'registered_users_m']:.0f}M+ users seed the enterprise funnel.",
                "orange",
            ),
            ("02", "Integrated platform", "Agent, workspace, cloud, and governance.", "blue"),
            (
                "03",
                "Improving margins",
                f"Lower model cost and operating scale produce {financials.loc[2031, 'gross_margin']:.0%} gross margin.",
                "coral",
            ),
            (
                "04",
                "Large market",
                f"${total_tam / 1000:.0f}B TAM with several revenue sources.",
                "green",
            ),
        ]
        for index, (num, title, body, color) in enumerate(thesis):
            x = 0.76 + index * 3.03
            add_rect(slide, x, 3.70, 2.70, 1.44, fill="charcoal", radius=True)
            add_text(
                slide,
                num,
                x + 0.20,
                3.93,
                0.38,
                0.22,
                font_size=8,
                color=color,
                bold=True,
                font_name=FONT_MONO,
            )
            add_text(
                slide,
                title,
                x + 0.20,
                4.27,
                2.25,
                0.28,
                font_size=10.5,
                color="white",
                bold=True,
            )
            add_text(
                slide,
                body,
                x + 0.20,
                4.65,
                2.25,
                0.34,
                font_size=7.4,
                color="gray_300",
            )
        add_rect(slide, 0.76, 5.66, 11.78, 0.81, fill="orange", radius=True)
        add_text(
            slide,
            "$1B Series E",
            1.05,
            5.87,
            2.30,
            0.36,
            font_size=18,
            color="white",
            bold=True,
        )
        add_text(
            slide,
            "funds model training, user growth, enterprise sales, international expansion, and talent",
            3.42,
            5.91,
            8.55,
            0.31,
            font_size=10.5,
            color="white",
            bold=True,
            align=PP_ALIGN.RIGHT,
        )


def build_presentation(
    model: ModelBundle,
    sources: dict[str, Source],
    slide_content: dict[str, Any],
    output_path: Path = PPTX_PATH,
) -> Path:
    return DeckBuilder(model, sources, slide_content).build(output_path)
