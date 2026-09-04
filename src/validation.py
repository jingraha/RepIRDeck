from __future__ import annotations

import math
from pathlib import Path
from typing import Any

from openpyxl import load_workbook
from pptx import Presentation

from .config import ACTIVE_SLIDE_IDS
from .model import ModelBundle
from .sources import Source


class ValidationError(RuntimeError):
    pass


def _assert(condition: bool, message: str) -> None:
    if not condition:
        raise ValidationError(message)


def validate_inputs(
    model: ModelBundle,
    sources: dict[str, Source],
    slide_content: dict[str, Any],
) -> None:
    slides = slide_content["slides"]
    source_map = slide_content["source_map"]

    _assert(len(slides) == 30, "The slide-content library must contain 30 source slides")
    _assert(set(slides) == set(range(1, 31)), "Slide numbers must be 1 through 30")
    _assert(
        set(source_map) == set(range(1, 31)),
        "Every slide must have a source-map entry",
    )

    missing_sources = {
        source_id
        for ids in source_map.values()
        for source_id in ids
        if source_id not in sources
    }
    _assert(not missing_sources, f"Unknown source IDs: {sorted(missing_sources)}")

    for source in sources.values():
        if source.evidence_class == "public":
            _assert(bool(source.url), f"Public source {source.id} must have a URL")

    transaction = model.transaction
    _assert(
        math.isclose(
            transaction["primary"] + transaction["secondary"],
            transaction["total_raise"],
        ),
        "Primary plus secondary must equal total raise",
    )
    _assert(
        math.isclose(
            transaction["post_money"],
            transaction["pre_money"] + transaction["primary"],
        ),
        "Post-money valuation must include primary capital only",
    )
    _assert(
        math.isclose(model.use_of_proceeds["amount"].sum(), transaction["primary"]),
        "Use of proceeds must equal primary capital",
    )

    financials = model.financials
    segment_total = financials[
        [
            "revenue_individual",
            "revenue_enterprise",
            "revenue_usage_and_deploy",
        ]
    ].sum(axis=1)
    _assert(
        (segment_total.round(6) == financials["revenue"].round(6)).all(),
        "Revenue segments do not reconcile",
    )
    _assert(
        financials["gross_margin"].between(-0.25, 0.90).all(),
        "Gross margin is outside a reasonable validation range",
    )

    scenario_years = set(range(2026, 2032))
    for name, scenario in model.scenarios.items():
        _assert(
            set(scenario.index) == scenario_years,
            f"{name} scenario must cover 2026 through 2031",
        )
        _assert(
            scenario["gross_margin"].between(0, 0.90).all(),
            f"{name} scenario gross margin is invalid",
        )

    _assert(
        model.scenarios["downside"].loc[2031, "revenue"]
        < model.scenarios["base"].loc[2031, "revenue"]
        < model.scenarios["upside"].loc[2031, "revenue"],
        "Scenario revenue ordering is invalid",
    )
    _assert(
        model.monthly_cash_flow["primary_financing"].sum()
        == transaction["primary"],
        "Monthly cash flow must include the primary financing once",
    )
    _assert(
        model.cash_summary["funding_need_date"] is not None,
        "The no-financing case must show a funding need",
    )
    _assert(
        model.cash_summary["minimum_base_cash"]
        >= model.cash_summary["minimum_cash"],
        "Base cash falls below the minimum operating cash level",
    )
    _assert(
        model.cash_summary["minimum_downside_cash"]
        >= model.cash_summary["minimum_cash"],
        "Downside cash falls below the minimum operating cash level",
    )
    _assert(
        model.balance_sheet["balance_check"].abs().max() < 0.000001,
        "Balance sheet does not balance",
    )
    _assert(
        model.cash_flow_statement.loc[2026, "ending_cash"]
        == model.monthly_cash_flow.loc["2026-12-31", "ending_cash"],
        "Monthly cash flow does not tie to annual cash flow",
    )


def validate_outputs(
    pptx_path: Path,
    xlsx_path: Path,
    expected_sheet_names: list[str],
) -> None:
    _assert(pptx_path.exists(), f"Missing presentation: {pptx_path}")
    _assert(xlsx_path.exists(), f"Missing workbook: {xlsx_path}")
    _assert(pptx_path.stat().st_size > 100_000, "Presentation appears incomplete")
    _assert(xlsx_path.stat().st_size > 25_000, "Workbook appears incomplete")

    presentation = Presentation(pptx_path)
    _assert(
        len(presentation.slides) == len(ACTIVE_SLIDE_IDS),
        f"Presentation must contain {len(ACTIVE_SLIDE_IDS)} slides",
    )
    for slide_number, slide in enumerate(presentation.slides, 1):
        for shape in slide.shapes:
            _assert(shape.left >= 0, f"Slide {slide_number} has a shape left of canvas")
            _assert(shape.top >= 0, f"Slide {slide_number} has a shape above canvas")
            _assert(
                shape.left + shape.width <= presentation.slide_width,
                f"Slide {slide_number} has a shape beyond the right edge",
            )
            _assert(
                shape.top + shape.height <= presentation.slide_height,
                f"Slide {slide_number} has a shape below the bottom edge",
            )

    workbook = load_workbook(xlsx_path, data_only=False, read_only=True)
    _assert(
        workbook.sheetnames == expected_sheet_names,
        "Workbook tab order or names are incorrect",
    )
    workbook.close()
