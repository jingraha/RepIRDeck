from __future__ import annotations

from collections.abc import Sequence
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .config import ASSET_DIR, COLORS, FONT_MONO, FONT_SANS


def rgb(color: str) -> RGBColor:
    return RGBColor.from_string(COLORS.get(color, color))


def set_background(slide: Any, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(
    slide: Any,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: float = 12,
    color: str = "ink",
    bold: bool = False,
    font_name: str = FONT_SANS,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0,
    italic: bool = False,
    all_caps: bool = False,
    line_spacing: float = 1.0,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text.upper() if all_caps else text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(
    slide: Any,
    runs: Sequence[dict[str, Any]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for spec in runs:
        run = paragraph.add_run()
        run.text = str(spec["text"])
        run.font.name = spec.get("font_name", FONT_SANS)
        run.font.size = Pt(spec.get("font_size", 12))
        run.font.bold = spec.get("bold", False)
        run.font.italic = spec.get("italic", False)
        run.font.color.rgb = rgb(spec.get("color", "ink"))
    return box


def add_rect(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = "white",
    line: str | None = None,
    radius: bool = False,
    transparency: int = 0,
) -> Any:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    if radius and shape.adjustments:
        shape.adjustments[0] = 0.12
    return shape


def add_circle(
    slide: Any,
    x: float,
    y: float,
    diameter: float,
    *,
    fill: str = "white",
    line: str | None = None,
    transparency: int = 0,
) -> Any:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x),
        Inches(y),
        Inches(diameter),
        Inches(diameter),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_line(
    slide: Any,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = "gray_300",
    width: float = 1,
    transparency: int = 0,
    dash: bool = False,
) -> Any:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    line.line.transparency = transparency
    if dash:
        line.line.dash_style = 4
    return line


def add_chevron(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = "orange",
) -> Any:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()
    return shape


def add_tag(
    slide: Any,
    text: str,
    x: float,
    y: float,
    *,
    fill: str = "gray_100",
    color: str = "gray_700",
    width: float | None = None,
) -> Any:
    width = width or max(0.55, 0.08 * len(text) + 0.22)
    add_rect(slide, x, y, width, 0.25, fill=fill, radius=True)
    add_text(
        slide,
        text,
        x,
        y + 0.01,
        width,
        0.20,
        font_size=7,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        all_caps=True,
    )
    return width


def add_metric_card(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    value: str,
    label: str,
    *,
    accent: str = "orange",
    evidence: str | None = None,
    fill: str = "white",
    dark: bool = False,
    value_size: float = 25,
    detail: str | None = None,
) -> None:
    add_rect(slide, x, y, w, h, fill=fill, line=None if dark else "gray_200", radius=True)
    add_rect(slide, x, y, 0.07, h, fill=accent, radius=True)
    value_color = "white" if dark else "ink"
    label_color = "gray_300" if dark else "gray_700"
    add_text(
        slide,
        value,
        x + 0.25,
        y + 0.18,
        w - 0.40,
        0.48,
        font_size=value_size,
        color=value_color,
        bold=True,
    )
    add_text(
        slide,
        label,
        x + 0.25,
        y + 0.73,
        w - 0.40,
        0.40,
        font_size=9,
        color=label_color,
        bold=True,
    )
    if detail:
        add_text(
            slide,
            detail,
            x + 0.25,
            y + h - 0.45,
            w - 0.40,
            0.30,
            font_size=7.5,
            color=label_color,
        )
    if evidence:
        add_tag(
            slide,
            evidence,
            x + w - max(0.68, 0.07 * len(evidence) + 0.22) - 0.12,
            y + 0.12,
            fill=accent if dark else "gray_100",
            color="white" if dark else "gray_700",
        )


def add_card_title(
    slide: Any,
    index: str,
    title: str,
    body: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    accent: str = "orange",
    fill: str = "white",
) -> None:
    add_rect(slide, x, y, w, h, fill=fill, line="gray_200", radius=True)
    add_text(
        slide,
        index,
        x + 0.20,
        y + 0.18,
        0.38,
        0.24,
        font_size=8,
        color=accent,
        bold=True,
        font_name=FONT_MONO,
    )
    add_text(
        slide,
        title,
        x + 0.20,
        y + 0.48,
        w - 0.40,
        0.36,
        font_size=13,
        color="ink",
        bold=True,
    )
    add_text(
        slide,
        body,
        x + 0.20,
        y + 0.88,
        w - 0.40,
        h - 1.02,
        font_size=9,
        color="gray_700",
        line_spacing=1.0,
    )


def add_replit_mark(
    slide: Any,
    x: float,
    y: float,
    *,
    scale: float = 1,
    wordmark: bool = True,
    color: str = "orange",
    wordmark_color: str = "ink",
) -> None:
    if wordmark:
        asset_name = (
            "replit_logo_on_dark.png"
            if wordmark_color in ("white", "cream")
            else "replit_logo_on_light.png"
        )
        asset_path = ASSET_DIR / "brand" / asset_name
        if asset_path.exists():
            slide.shapes.add_picture(
                str(asset_path),
                Inches(x),
                Inches(y),
                width=Inches(2.05 * scale),
            )
            return
    else:
        asset_path = ASSET_DIR / "brand" / "replit_symbol_orange.png"
        if asset_path.exists():
            slide.shapes.add_picture(
                str(asset_path),
                Inches(x),
                Inches(y),
                width=Inches(0.55 * scale),
            )
            return

    block_h = 0.13 * scale
    block_w = 0.55 * scale
    gap = 0.09 * scale
    offsets = [0.0, 0.12 * scale, 0.0]
    for index, offset in enumerate(offsets):
        add_rect(
            slide,
            x + offset,
            y + index * (block_h + gap),
            block_w - offset,
            block_h,
            fill=color,
            radius=True,
        )
    if wordmark:
        add_text(
            slide,
            "replit",
            x + 0.72 * scale,
            y - 0.02 * scale,
            1.4 * scale,
            0.62 * scale,
            font_size=21 * scale,
            color=wordmark_color,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )


def add_native_line_chart(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    labels: Sequence[str],
    series: Sequence[dict[str, Any]],
    *,
    min_value: float | None = None,
    max_value: float | None = None,
    y_format: str = "{:.0f}",
    show_legend: bool = True,
    show_end_labels: bool = True,
    gridlines: int = 4,
    dark: bool = False,
) -> None:
    all_values = [float(value) for item in series for value in item["values"]]
    y_min = min_value if min_value is not None else min(all_values)
    y_max = max_value if max_value is not None else max(all_values)
    if y_max == y_min:
        y_max += 1
    axis_color = "gray_500" if dark else "gray_300"
    label_color = "gray_300" if dark else "gray_700"
    plot_left = x + 0.48
    plot_top = y + (0.35 if show_legend else 0.08)
    plot_width = w - 0.72
    plot_height = h - (0.72 if show_legend else 0.45)

    for line_index in range(gridlines + 1):
        ratio = line_index / gridlines
        y_pos = plot_top + plot_height - ratio * plot_height
        add_line(
            slide,
            plot_left,
            y_pos,
            plot_left + plot_width,
            y_pos,
            color=axis_color,
            width=0.6,
            transparency=45,
        )
        value = y_min + ratio * (y_max - y_min)
        add_text(
            slide,
            y_format.format(value),
            x,
            y_pos - 0.10,
            0.40,
            0.20,
            font_size=6.5,
            color=label_color,
            align=PP_ALIGN.RIGHT,
        )

    if len(labels) == 1:
        x_positions = [plot_left + plot_width / 2]
    else:
        x_positions = [
            plot_left + index * plot_width / (len(labels) - 1)
            for index in range(len(labels))
        ]
    for x_pos, label in zip(x_positions, labels, strict=True):
        add_text(
            slide,
            label,
            x_pos - 0.28,
            plot_top + plot_height + 0.10,
            0.56,
            0.20,
            font_size=6.5,
            color=label_color,
            align=PP_ALIGN.CENTER,
        )

    legend_x = plot_left
    for item in series:
        values = [float(value) for value in item["values"]]
        color = item.get("color", "orange")
        points = []
        for x_pos, value in zip(x_positions, values, strict=True):
            ratio = (value - y_min) / (y_max - y_min)
            y_pos = plot_top + plot_height - ratio * plot_height
            points.append((x_pos, y_pos))
        for point_index in range(len(points) - 1):
            add_line(
                slide,
                points[point_index][0],
                points[point_index][1],
                points[point_index + 1][0],
                points[point_index + 1][1],
                color=color,
                width=item.get("width", 2.3),
            )
        for x_pos, y_pos in points:
            add_circle(
                slide,
                x_pos - 0.055,
                y_pos - 0.055,
                0.11,
                fill=color,
                line="white" if not dark else "ink",
            )
        if show_end_labels:
            add_text(
                slide,
                item.get("end_label", y_format.format(values[-1])),
                points[-1][0] - 0.45,
                points[-1][1] - 0.36,
                0.9,
                0.22,
                font_size=7.5,
                color=color,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
        if show_legend:
            add_rect(
                slide,
                legend_x,
                y,
                0.18,
                0.06,
                fill=color,
                radius=True,
            )
            add_text(
                slide,
                item["name"],
                legend_x + 0.25,
                y - 0.07,
                1.25,
                0.20,
                font_size=7,
                color=label_color,
            )
            legend_x += 1.55


def add_stacked_columns(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    labels: Sequence[str],
    series: Sequence[dict[str, Any]],
    *,
    max_value: float | None = None,
    value_format: str = "${:,.0f}",
    dark: bool = False,
) -> None:
    totals = [
        sum(float(item["values"][index]) for item in series)
        for index in range(len(labels))
    ]
    y_max = max_value or max(totals) * 1.10
    label_color = "gray_300" if dark else "gray_700"
    baseline = y + h - 0.40
    plot_height = h - 0.75
    column_slot = w / len(labels)
    column_width = min(0.62, column_slot * 0.55)
    add_line(slide, x, baseline, x + w, baseline, color="gray_300", width=0.7)

    for index, label in enumerate(labels):
        left = x + index * column_slot + (column_slot - column_width) / 2
        running_height = 0.0
        for item in series:
            value = float(item["values"][index])
            segment_height = value / y_max * plot_height
            add_rect(
                slide,
                left,
                baseline - running_height - segment_height,
                column_width,
                segment_height,
                fill=item.get("color", "orange"),
                radius=segment_height > 0.22,
            )
            running_height += segment_height
        add_text(
            slide,
            value_format.format(totals[index]),
            left - 0.18,
            baseline - running_height - 0.28,
            column_width + 0.36,
            0.22,
            font_size=7,
            color="white" if dark else "ink",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            label,
            left - 0.12,
            baseline + 0.10,
            column_width + 0.24,
            0.20,
            font_size=6.5,
            color=label_color,
            align=PP_ALIGN.CENTER,
        )


def add_horizontal_bars(
    slide: Any,
    x: float,
    y: float,
    w: float,
    labels: Sequence[str],
    values: Sequence[float],
    colors: Sequence[str],
    *,
    max_value: float | None = None,
    value_format: str = "{:,.0f}",
    bar_height: float = 0.30,
    row_height: float = 0.62,
    dark: bool = False,
) -> None:
    max_value = max_value or max(values)
    text_color = "gray_300" if dark else "gray_700"
    for index, (label, value, color) in enumerate(zip(labels, values, colors, strict=True)):
        row_y = y + index * row_height
        add_text(
            slide,
            label,
            x,
            row_y,
            1.75,
            0.28,
            font_size=8,
            color=text_color,
            bold=True,
        )
        track_x = x + 1.85
        track_w = w - 2.75
        add_rect(
            slide,
            track_x,
            row_y + 0.02,
            track_w,
            bar_height,
            fill="gray_100" if not dark else "charcoal",
            radius=True,
        )
        add_rect(
            slide,
            track_x,
            row_y + 0.02,
            max(0.04, track_w * value / max_value),
            bar_height,
            fill=color,
            radius=True,
        )
        add_text(
            slide,
            value_format.format(value),
            x + w - 0.78,
            row_y - 0.02,
            0.78,
            0.30,
            font_size=8.5,
            color="white" if dark else "ink",
            bold=True,
            align=PP_ALIGN.RIGHT,
        )


def add_matrix_dot(
    slide: Any,
    x: float,
    y: float,
    level: int,
    *,
    active_color: str = "orange",
    inactive_color: str = "gray_200",
) -> None:
    for index in range(3):
        add_circle(
            slide,
            x + index * 0.15,
            y,
            0.095,
            fill=active_color if index < level else inactive_color,
        )


def add_bullet_list(
    slide: Any,
    items: Sequence[str],
    x: float,
    y: float,
    w: float,
    *,
    font_size: float = 9,
    color: str = "gray_700",
    bullet_color: str = "orange",
    row_height: float = 0.48,
) -> None:
    for index, item in enumerate(items):
        row_y = y + index * row_height
        add_circle(slide, x, row_y + 0.08, 0.08, fill=bullet_color)
        add_text(
            slide,
            item,
            x + 0.18,
            row_y,
            w - 0.18,
            row_height,
            font_size=font_size,
            color=color,
        )
