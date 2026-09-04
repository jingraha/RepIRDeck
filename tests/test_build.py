from __future__ import annotations

from openpyxl import load_workbook
from pptx import Presentation

from src.build import build_all
from src.xlsx_builder import SHEET_NAMES


def test_build_outputs(tmp_path) -> None:
    pptx_path, xlsx_path = build_all(tmp_path)
    presentation = Presentation(pptx_path)
    assert len(presentation.slides) == 24
    slide_text = {
        index + 1: "\n".join(
            shape.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        for index, slide in enumerate(presentation.slides)
    }
    assert "$20.8B" in slide_text[2]
    assert "52%" in slide_text[16]
    assert "Nov-26" in slide_text[21]
    assert "H2 2030" in slide_text[21]
    assert "Enterprise and runtime drive the 2026 to 2031 plan" in slide_text[16]
    assert "Enterprise accounts expand after the first deployment" in slide_text[17]
    assert "Lower model cost lifts gross margin" in slide_text[18]
    all_slide_text = "\n".join(slide_text.values())
    assert "UKG reports 4x more customer feedback before engineering starts" not in all_slide_text
    assert "Replit charges for access, agent work, and running apps" not in all_slide_text
    assert "$880M of pipeline supports the 2027 plan" not in all_slide_text
    assert "Usage data, runtime, and controls make Replit harder to replace" not in all_slide_text
    assert "Revenue grows across three distinct sources" not in all_slide_text
    assert "Subscriptions provide a base. Usage adds growth." not in all_slide_text
    assert "Residual" not in slide_text[23]

    workbook = load_workbook(xlsx_path, read_only=True, data_only=False)
    assert workbook.sheetnames == SHEET_NAMES
    assert workbook.sheetnames[:5] == [
        "Assumptions",
        "Visuals",
        "P&L",
        "Cash Flow",
        "Balance Sheet",
    ]
    assert workbook["Cash Flow"]["A1"].value == "Cash Flow"
    assert workbook["Cash Flow"]["B5"].value == "='P&L'!B15"
    assert workbook["Balance Sheet"]["B5"].value == "='Cash Flow'!B16"
    reconciliation_slides = {
        row[0]
        for row in workbook["Slide Reconciliation"].iter_rows(
            min_row=5,
            values_only=True,
        )
        if row[0]
    }
    assert max(reconciliation_slides) == 24
    assert all(1 <= slide <= 24 for slide in reconciliation_slides)
    assert workbook["QA Checks"]["D5"].value == "PASS"
    workbook.close()
